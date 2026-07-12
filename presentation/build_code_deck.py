#!/usr/bin/env python3
"""Builds lab_meeting_full.pptx: a full code-walkthrough deck matching the
style of "260713 lab meeting.pptx" (white background, black Arial text,
VS-Code-dark code panels in Consolas, Courier New console-style results).

For every meaningful analysis step in bioinform25/project1, shows: the full
R source code that produced a result, a Korean explanation of what the code
does, and the resulting figure/table with a Korean explanation of what it
shows. English titles throughout.
"""

import os
import re
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

REPO = os.path.join(os.path.dirname(__file__), "..")
FIG = os.path.join(REPO, "figures")
RES = os.path.join(REPO, "results")
EXT_FIG = os.path.join(REPO, "external_validation", "figures")
EXT_RES = os.path.join(REPO, "external_validation", "results")
SCRIPTS = os.path.join(REPO, "scripts")
EXT_SCRIPTS = os.path.join(REPO, "external_validation", "scripts")
OUT_PATH = os.path.join(os.path.dirname(__file__), "lab_meeting_full.pptx")

# ------------------------------------------------------------------ style ---
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG = RGBColor(0x1E, 0x1E, 0x1E)
C_COMMENT = RGBColor(0x6A, 0x99, 0x55)
C_FUNC = RGBColor(0xDC, 0xDC, 0xAA)
C_STRING = RGBColor(0xCE, 0x91, 0x78)
C_KEYWORD = RGBColor(0x56, 0x9C, 0xD6)
C_PLAIN = RGBColor(0xD4, 0xD4, 0xD4)

FONT_UI = "Arial"
FONT_CODE = "Consolas"
FONT_RESULT = "Courier New"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def new_slide():
    s = prs.slides.add_slide(BLANK)
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    return s


def add_textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def set_run(run, text, size, color, bold=False, italic=False, font=FONT_UI):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color


def slide_title(slide, text):
    tb, tf = add_textbox(slide, Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.55))
    p = tf.paragraphs[0]
    r = p.add_run()
    set_run(r, text, 20, BLACK, bold=True)


def page_number(slide, n):
    tb, tf = add_textbox(slide, Inches(12.7), Inches(7.15), Inches(0.5), Inches(0.3))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    set_run(r, str(n), 10, RGBColor(0x80, 0x80, 0x80))


def add_explanation(slide, text, left=Inches(7.95), top=Inches(1.05), width=Inches(4.98), height=Inches(1.4), size=11.5):
    tb, tf = add_textbox(slide, left, top, width, height)
    p = tf.paragraphs[0]
    p.line_spacing = 1.15
    r = p.add_run()
    set_run(r, text, size, BLACK)
    return tb


def add_result_image(slide, path, left=Inches(7.95), top=Inches(2.58), max_w=Inches(4.98), max_h=Inches(3.91)):
    with Image.open(path) as im:
        iw, ih = im.size
    ar = iw / ih
    box_ar = max_w / max_h
    if ar > box_ar:
        w = max_w
        h = int(max_w / ar)
    else:
        h = max_h
        w = int(max_h * ar)
    x = left + (max_w - w) // 2
    y = top + (max_h - h) // 2
    slide.shapes.add_picture(path, x, y, width=w, height=h)


def add_result_console(slide, lines, left=Inches(7.95), top=Inches(2.58), width=Inches(4.98), height=Inches(3.91), size=10):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(0)
        r = p.add_run()
        set_run(r, line, size, BLACK, font=FONT_RESULT)
    return shp


def add_result_table(slide, headers, rows, left=Inches(7.95), top=Inches(2.58), width=Inches(4.98), height=Inches(3.91), size=9.5):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    gshape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gshape.table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h if h else " "
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.runs[0]
        set_run(run, h if h else " ", size, WHITE, bold=True, font=FONT_UI)
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x40, 0x40, 0x40)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for r_i, row in enumerate(rows, start=1):
        for c_i, val in enumerate(row):
            cell = table.cell(r_i, c_i)
            cell.text = str(val) if val != "" else " "
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.runs[0]
            set_run(run, str(val) if val != "" else " ", size, BLACK, font=FONT_UI)
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if r_i % 2 else RGBColor(0xF2, 0xF2, 0xF2)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return gshape


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ============================================================ R HIGHLIGHT ==
_KEYWORDS = (
    "TRUE|FALSE|NA|NA_character_|NA_real_|NA_integer_|NULL|Inf|NaN|"
    "if|else|for|function|return|break|next|in|repeat|while"
)
_TOKEN_RE = re.compile(
    r'(?P<string>"(?:[^"\\]|\\.)*"|\'(?:[^\'\\]|\\.)*\')'
    r'|(?P<comment>#.*$)'
    r'|(?P<keyword>\b(?:' + _KEYWORDS + r')\b)'
    r'|(?P<func>[A-Za-z_.][A-Za-z0-9_.]*(?=\s*\())'
    r'|(?P<other>.)'
)


def highlight_r_line(line):
    """Returns list of (text, color) runs for one line of R code."""
    segs = []
    for m in _TOKEN_RE.finditer(line):
        if m.lastgroup == "string":
            color = C_STRING
        elif m.lastgroup == "comment":
            color = C_COMMENT
        elif m.lastgroup == "keyword":
            color = C_KEYWORD
        elif m.lastgroup == "func":
            color = C_FUNC
        else:
            color = C_PLAIN
        text = m.group(0)
        if segs and segs[-1][1] == color:
            segs[-1] = (segs[-1][0] + text, color)
        else:
            segs.append((text, color))
    if not segs:
        segs = [(" ", C_PLAIN)]
    return segs


CODE_FONT_SIZE = 6.5
CODE_MAX_LINES = 50


def paginate_code(lines, max_lines=CODE_MAX_LINES):
    """Split a list of source lines into pages, preferring blank-line breaks
    near the cap so a page doesn't end mid logical-block."""
    if len(lines) <= max_lines:
        return [lines]
    pages = []
    remaining = lines
    while len(remaining) > max_lines:
        cut = max_lines
        for back in range(0, 8):
            idx = max_lines - back
            if idx > 0 and remaining[idx - 1].strip() == "":
                cut = idx
                break
        pages.append(remaining[:cut])
        remaining = remaining[cut:]
    if remaining:
        pages.append(remaining)
    return pages


def add_code_block(slide, code_lines, left=Inches(0.4), top=Inches(1.05),
                    width=Inches(7.35), height=Inches(5.7)):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = CODE_BG
    shp.line.fill.background()
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = 0.02
    except Exception:
        pass
    tf = shp.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    for i, line in enumerate(code_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(0)
        p.line_spacing = 1.05
        for text, color in highlight_r_line(line.rstrip("\n")):
            r = p.add_run()
            set_run(r, text, CODE_FONT_SIZE, color, font=FONT_CODE)
    return shp


def read_script(path):
    with open(path, encoding="utf-8") as f:
        return f.read().splitlines()


import csv as _csv


def read_csv_preview(path, max_rows=8, max_cols=7):
    with open(path, encoding="utf-8") as f:
        rows = list(_csv.reader(f))
    headers = [h.strip('"') for h in rows[0]][:max_cols]
    body = [[c.strip('"') for c in r][:max_cols] for r in rows[1:1 + max_rows]]
    truncated_cols = len(rows[0]) > max_cols
    truncated_rows = len(rows) - 1 > max_rows
    if truncated_cols:
        headers.append("...")
        body = [r + ["..."] for r in body]
    return headers, body, truncated_rows, len(rows) - 1


# ==================================================================== SLIDES ==
_slide_counter = [0]


def title_slide(title, subtitle, footer):
    s = new_slide()
    tb, tf = add_textbox(s, Inches(1.0), Inches(2.4), Inches(11.3), Inches(3.2))
    p = tf.paragraphs[0]
    r = p.add_run(); set_run(r, title, 30, BLACK, bold=True)
    p2 = tf.add_paragraph(); p2.space_before = Pt(16)
    r2 = p2.add_run(); set_run(r2, subtitle, 18, BLACK)
    p3 = tf.add_paragraph(); p3.space_before = Pt(10)
    r3 = p3.add_run(); set_run(r3, footer, 13, BLACK)
    return s


def flow_overview_slide(title, lines):
    s = new_slide()
    slide_title(s, title)
    tb, tf = add_textbox(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(5.9))
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        r = p.add_run(); set_run(r, "- " + line, 15, BLACK)
    _slide_counter[0] += 1
    page_number(s, _slide_counter[0])
    return s


def section_divider(title, subtitle):
    s = new_slide()
    tb, tf = add_textbox(s, Inches(0.8), Inches(3.0), Inches(11.7), Inches(2.0))
    p = tf.paragraphs[0]
    r = p.add_run(); set_run(r, title, 32, BLACK, bold=True)
    p2 = tf.add_paragraph(); p2.space_before = Pt(10)
    r2 = p2.add_run(); set_run(r2, subtitle, 16, BLACK)
    _slide_counter[0] += 1
    page_number(s, _slide_counter[0])
    return s


def code_result_slide(title, code_lines, explain_kr, result_fn, result_explain_kr,
                       part_label=None):
    """One slide: title, dark code panel (left), Korean explanation (right-top),
    result image/table/console (right-bottom) + its Korean explanation appended."""
    s = new_slide()
    full_title = title if not part_label else f"{title} ({part_label})"
    slide_title(s, full_title)
    add_code_block(s, code_lines)
    explain_h = Inches(1.5)
    add_explanation(s, explain_kr, top=Inches(1.05), height=explain_h)
    if result_fn is not None:
        result_fn(s)
    if result_explain_kr:
        add_explanation(s, result_explain_kr, left=Inches(7.95), top=Inches(6.55),
                         width=Inches(4.98), height=Inches(0.85), size=10.5)
    _slide_counter[0] += 1
    page_number(s, _slide_counter[0])
    return s


def emit_step(title, script_path, explain_kr, result_fn=None, result_explain_kr=None):
    """Reads script_path, paginates its full source, and emits one slide per
    page. Only the LAST page carries the result (image/table/console) and its
    explanation; earlier pages carry only code + the same code explanation."""
    lines = read_script(script_path)
    pages = paginate_code(lines)
    n = len(pages)
    for i, page_lines in enumerate(pages):
        part_label = f"{i+1}/{n}" if n > 1 else None
        is_last = (i == n - 1)
        code_result_slide(
            title, page_lines, explain_kr,
            result_fn if is_last else None,
            result_explain_kr if is_last else None,
            part_label=part_label,
        )


def summary_slide(title, sections):
    """sections: list of (heading, [lines])"""
    s = new_slide()
    slide_title(s, title)
    tb, tf = add_textbox(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(5.9))
    first = True
    for heading, lines in sections:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(0 if first else 14)
        r = p.add_run(); set_run(r, heading, 16, BLACK, bold=True)
        for line in lines:
            p2 = tf.add_paragraph()
            p2.space_after = Pt(8)
            r2 = p2.add_run(); set_run(r2, "- " + line, 14, BLACK)
    _slide_counter[0] += 1
    page_number(s, _slide_counter[0])
    return s


def closing_slide(title, subtitle):
    s = new_slide()
    tb, tf = add_textbox(s, Inches(1.0), Inches(3.1), Inches(11.3), Inches(1.6))
    p = tf.paragraphs[0]
    r = p.add_run(); set_run(r, title, 34, BLACK, bold=True)
    p2 = tf.add_paragraph(); p2.space_before = Pt(12)
    r2 = p2.add_run(); set_run(r2, subtitle, 16, BLACK)
    return s

# =================================================================== MAIN ==
def build():
    # ---------------------------------------------------------------- title
    title_slide(
        "S100A8+ Macrophage & MASLD: Full Code Walkthrough",
        "scRNA-Informed Deconvolution + Origin/Trigger/Fate Follow-up (GSE285614 + GSE285615)",
        "bioinform25/project1  |  Source paper: Guan et al., J Clin Invest 2025;135(21):e190635",
    )

    # ----------------------------------------------------------- flow slide
    flow_overview_slide("Overall Study Flow", [
        "원 논문(Guan et al., JCI 2025): 지방세포 사멸 -> S100A8+ 대식세포 -> CCN3 억제 -> CD36 상승 -> 간세포 지방축적 -> MASLD 진행 (CD36/CCN3 단일 유전자 경로를 bulk DEG로 추적)",
        "우리 질문: 같은 마우스 코호트의 bulk(GSE285614, 34샘플, 2개 genotype axis)와 scRNA(GSE285615, WT/TG 각 1개)를 결합해, 세포 구성 자체가 genotype에 따라 바뀌는지 직접 물을 수 있지 않을까?",
        "1부 (scripts/01-11): scRNA 처리 -> marker-based deconvolution(BisqueRNA) -> axis-1(WT/TG)로 검증 -> axis-2(MKO)에 적용 -> 지방생성 유전자 상관분석 -> 통계적 엄밀성 재점검",
        "2부 (external_validation/01-05): 독립된 외부 scRNA 데이터셋 2개(GSE156057 CD45+ 정렬, GSE166504 whole-liver)로 S100A8hi_macrophage population 자체의 재현성 검증",
        "3부 (scripts/13-15, external_validation/06): 원 논문이 명시적으로 미해결로 남긴 3가지 질문(기원/유발인자/운명)을 같은 데이터로 후속 분석",
        "핵심 발견: 지방생성 유전자와 상관된 것은 S100A8hi 서브셋이 아니라 일반 대식세포/Kupffer 총량 -- 그리고 그 \"Kupffer\" 안에 숨어있던 이질성을 후속 분석에서 발견",
    ])

    # =============================================== PART 1: PRIMARY ==
    section_divider("Part 1 — Primary Analysis", "GSE285614 (bulk RNA-seq, 34 samples) + GSE285615 (scRNA-seq, WT vs TG)")

    emit_step(
        "01_scRNA_qc.R — Load & QC-filter GSE285615 scRNA-seq",
        os.path.join(SCRIPTS, "01_scRNA_qc.R"),
        "GSE285615의 WT_HFD, TG_HFD 두 샘플 h5 파일(GSM8705804, GSM8705805)에서, 유효한 세포만 "
        "남기기 위해 nFeature_RNA(200~6000), nCount_RNA(<30000), percent.mt(<25%) 범위로 QC "
        "필터링을 했다. 간세포는 미토콘드리아 비중이 자연히 높아 mt 컷오프를 25%까지 넉넉히 잡았다.",
        result_fn=lambda s: add_result_image(s, os.path.join(FIG, "01_qc_violin_prefilter.png")),
        result_explain_kr="필터링 전/후 세포 수와 QC 지표(nFeature/nCount/percent.mt) 분포를 violin plot으로 확인해 두 샘플 모두 정상 범위임을 검증했다.",
    )

    emit_step(
        "02_scRNA_cluster_annotate.R — Cluster & annotate major lineages",
        os.path.join(SCRIPTS, "02_scRNA_cluster_annotate.R"),
        "QC를 통과한 WT/TG Seurat 객체를 병합한 뒤(샘플=genotype이라 배치보정 없이 단순 병합), "
        "PCA(30개 성분) -> FindNeighbors -> FindClusters(resolution=0.8) -> UMAP 범위에서 "
        "비지도 클러스터링을 했다. 간세포/Kupffer/호중구/단핵구/LSEC/성상세포 등 22개 canonical "
        "마커 유전자로 각 클러스터의 대분류 계통을 확인했다.",
        result_fn=lambda s: add_result_image(s, os.path.join(FIG, "02_umap_clusters.png")),
        result_explain_kr="28개 클러스터가 나왔고, FindAllMarkers로 클러스터별 top5 마커 유전자를 뽑아 이후 세포타입 라벨링(03단계)의 근거로 저장했다.",
    )

    emit_step(
        "03_scRNA_annotate_finalize.R — Identify S100A8hi_macrophage vs Neutrophil",
        os.path.join(SCRIPTS, "03_scRNA_annotate_finalize.R"),
        "02단계 클러스터별 top5 마커(results/02_cluster_top5_markers.csv)와 S100a8/S100a9/Ly6g/"
        "Csf1r/Adgre1/Clec4f/Trem2/Gpnmb 등 발현 비율을 클러스터별로 직접 대조하는 범위에서, "
        "28개 클러스터를 12개 세포타입으로 매핑했다. 핵심적으로 S100a8/S100a9 양성 세포가 "
        "단핵구/대식세포 계통(클러스터12)과 성숙 호중구(클러스터23) 두 개로 분리됨을 확인했다.",
        result_fn=lambda s: add_result_image(s, os.path.join(FIG, "03_confirm_s100a8_identity.png")),
        result_explain_kr="S100A8hi_macrophage(Ly6g 음성)와 Neutrophil(Ly6g/Camp/Ltf/Ngp 양성)이 "
                           "명확히 분리됨을 dot plot으로 검증 -- 이 둘을 구분하지 않으면 핵심 population이 호중구에 오염된다.",
    )

    emit_step(
        "04_scRNA_ground_truth_props.R — scRNA-observed cell-type proportions (ground truth)",
        os.path.join(SCRIPTS, "04_scRNA_ground_truth_props.R"),
        "03단계에서 라벨링된 scrna_annotated.rds에서, WT/TG 각 샘플별 세포타입 구성비를 직접 "
        "세는 범위에서 genotype x cell_type proportion 표를 만들고, TG/WT 배율을 계산했다. "
        "이건 이후 07단계에서 bulk deconvolution 추정치와 대조할 \"ground truth\"로 쓰인다.",
        result_fn=lambda s: add_result_image(s, os.path.join(FIG, "04_ground_truth_barplot.png")),
        result_explain_kr="scRNA 기준 S100A8hi_macrophage는 TG가 WT보다 1.46배 많았다 -- "
                           "다만 이건 genotype당 n=1이라, 07단계에서 이 방향이 실제로 신뢰할 수 있는지 재검증하게 된다.",
    )

    emit_step(
        "05_bulk_load_merge.R — Parse & merge the two bulk RNA-seq axes",
        os.path.join(SCRIPTS, "05_bulk_load_merge.R"),
        "GSE285614의 두 xlsx 파일(HFD_TG_vs_WT.xlsx, MKO_HFD_vs_WT_HFD.xlsx)에서, 저자가 함께 "
        "올려둔 시료별 raw count 컬럼만 추출해 유전자 심볼 기준으로 합치는 범위에서, axis1(WT vs "
        "Bcl2TG)과 axis2(WT vs MKO) 34개 샘플을 하나의 count matrix로 병합했다. 두 axis는 서로 "
        "다른 시퀀싱 배치이므로 genotype/axis를 분리해서 기록하고 절대 풀링하지 않았다.",
        result_fn=lambda s: add_result_table(
            s, ["axis", "genotype", "diet", "n"],
            [["adipocyte_death", "WT", "Chow", "3"], ["adipocyte_death", "Bcl2TG", "Chow", "3"],
             ["adipocyte_death", "WT", "HFD", "5"], ["adipocyte_death", "Bcl2TG", "HFD", "5"],
             ["macrophage_KO", "WT", "Chow", "4"], ["macrophage_KO", "MKO", "Chow", "5"],
             ["macrophage_KO", "WT", "HFD", "5"], ["macrophage_KO", "MKO", "HFD", "4"]]),
        result_explain_kr="34개 샘플의 최종 axis x genotype x diet 구성. 12,505개 공통 유전자로 "
                           "두 xlsx의 유전자셋을 교집합해 하나의 count matrix(results/bulk_counts.rds)로 저장했다.",
    )

    emit_step(
        "06_deconvolution.R — Marker-based deconvolution (BisqueRNA)",
        os.path.join(SCRIPTS, "06_deconvolution.R"),
        "scRNA 참조가 genotype당 n=1이라 subject-variance가 필요한 MuSiC/Bisque "
        "ReferenceBasedDecomposition은 부적절하다고 판단해, FindAllMarkers로 뽑은 세포타입별 "
        "마커 유전자 집합을 bulk의 log2(CPM+1) 발현값 범위에 투영하는 "
        "BisqueRNA::MarkerBasedDecomposition을 사용해 34개 샘플 전체의 상대적 세포타입 존재비 "
        "점수를 산출했다.",
        result_fn=lambda s: add_result_console(s, [
            "12 cell types, n markers used per type:",
            "",
            "  B_cell            5   Mast_cell          9",
            "  Cholangiocyte    14   Neutrophil        25",
            "  Dendritic_cell    5   NK_cell           18",
            "  Endothelial       5   S100A8hi_macroph  18",
            "  Hepatocyte        6   Stellate_Fibrobl.  9",
            "  Macroph_Kupffer   7   T_cell             6",
            "",
            "-> results/06_deconvolution_props.csv",
            "   (34 samples x 12 cell types)",
        ]),
        result_explain_kr="각 세포타입의 마커 패널 크기(5~25개 유전자). 결과값은 세포타입 간 "
                           "절대비율이 아니라 상대적 scale의 PCA score라, 같은 세포타입을 시료 간 비교하는 데만 유효하다.",
    )

    emit_step(
        "07_validate_axis1.R — Validate: does deconvolution match a trustworthy ground truth?",
        os.path.join(SCRIPTS, "07_validate_axis1.R"),
        "04단계 scRNA ground truth(n=1이라 저신뢰)와 06단계 deconvolution 결과가 반대 방향으로 "
        "나온 것을 발견하고, GSE285614 axis1 xlsx에 있는 원저자의 정식 DESeq2 결과(n=5 vs n=5, "
        "고신뢰로) 범위에서 S100a8/Cd68/Cd36 등 핵심 마커 유전자의 log2FoldChange 방향을 직접 "
        "대조해 어느 쪽이 맞는 방향인지 재확인했다.",
        result_fn=lambda s: add_result_image(s, os.path.join(FIG, "07_validation_S100A8_boxplot.png")),
        result_explain_kr="원저자 DESeq2(n=5 vs n=5)는 S100a8/Cd68/Cd36 모두 WT>TG 방향 -- "
                           "이는 06단계 deconvolution과 일치하고, n=1 scRNA ground truth와는 반대였다. "
                           "n=1 쌍이 개체차로 우연히 반대 방향이 나온 것으로 결론짓고, deconvolution을 신뢰하기로 했다.",
    )

    emit_step(
        "08_apply_axis2_MKO.R — Apply the validated method to axis-2 (MKO, no scRNA)",
        os.path.join(SCRIPTS, "08_apply_axis2_MKO.R"),
        "axis2(대식세포 특이적 S100a8 KO, MKO)는 scRNA 자체가 없으므로, 07단계에서 검증된 "
        "deconvolution 결과를 그대로 적용하는 범위에서 WT vs MKO를 chow/HFD 각각에서 "
        "Welch t-test로 비교하고 BH로 다중비교 보정을 했다.",
        result_fn=lambda s: add_result_image(s, os.path.join(FIG, "08_MKO_key_celltypes.png")),
        result_explain_kr="S100A8hi_macrophage, Macrophage_Kupffer 모두 WT vs MKO 유의차 없음 -- "
                           "정직하게 null 결과로 보고했고, 11단계에서 이게 \"효과 없음\"이 아니라 \"표본크기 부족\"임을 power 분석으로 재확인했다.",
    )

    emit_step(
        "09_correlate_lipid_genes.R — Correlate cell-type abundance with lipogenic genes",
        os.path.join(SCRIPTS, "09_correlate_lipid_genes.R"),
        "원 논문의 핵심 기전(S100A8+ macrophage -> CCN3 억제 -> CD36 상승 -> 지방축적)이 우리 "
        "deconvolution 결과와도 이어지는지 확인하기 위해, 34개 bulk 시료의 log2(CPM+1) 발현값 "
        "범위에서 S100A8hi_macrophage/Macrophage_Kupffer 존재비 점수와 Cd36/Ccn3/Fasn/Plin2/"
        "Pparg/Scd1 발현의 Spearman 상관을 axis/diet 서브셋별로 계산했다.",
        result_fn=lambda s: add_result_image(s, os.path.join(FIG, "09_correlation_S100A8_vs_Cd36.png")),
        result_explain_kr="놀랍게도 S100A8hi_macrophage 자체는 Cd36과 거의 무관(r=0.11)했다 -- "
                           "이게 바로 \"S100A8hi 서브셋이 아니라 일반 대식세포가 지방생성과 연결된다\"는 핵심 발견의 첫 단서였다.",
    )

    emit_step(
        "11_statistical_rigor_addendum.R — Statistical rigor re-audit before writing the manuscript",
        os.path.join(SCRIPTS, "11_statistical_rigor_addendum.R"),
        "원고 작성 전 자체 재점검에서 발견한 갭들을 고치는 범위에서: (1) 09단계 지질유전자 "
        "상관분석(30개 검정)에 BH-FDR 보정을 추가했고, (2) 07/08단계의 소표본(n=4-5) 비교에 "
        "Wilcoxon(비모수)을 1차 지표로, Cohen's d와 95% CI를 병기했고, (3) 08단계 null 결과가 "
        "\"효과 없음\"인지 \"표본이 작아 못 봄\"인지 구분하기 위해 80% power 기준 최소 검출가능 "
        "Cohen's d(MDE)를 계산했다.",
        result_fn=lambda s: add_result_console(s, [
            "80% power, min. detectable",
            "Cohen's d (alpha=0.05):",
            "",
            "   n=4      n=5",
            "   2.38     2.02",
            "",
            "-> MKO axis2 (n=4-5) only detects",
            "   d >= 2.0-2.4 -> null result is",
            "   \"underpowered\", not \"no effect\"",
        ]),
        result_explain_kr="이 재점검 덕분에 08단계의 null 결과를 과도하게 해석하지 않고, "
                           "\"이 표본크기로는 큰 효과만 검출 가능하다\"는 정직하고 정확한 결론으로 다듬을 수 있었다.",
    )

    # =============================================== PART 2: EXTERNAL ==
    section_divider("Part 2 — External Validation", "GSE156057 (Remmerie, CD45+-sorted) + GSE166504 (Su, whole-liver)")

    emit_step(
        "ext/01_gse166504_extract_myeloid.R — Extract myeloid cells from GSE166504",
        os.path.join(EXT_SCRIPTS, "01_gse166504_extract_myeloid.R"),
        "GSE166504(마우스 간 NPC scRNA, chow/15주-NAFL/30주-NASH/34주, 진짜 생물학적 반복 "
        "n=3-8/그룹)의 전체 raw count 파일은 너무 커서, 저자가 제공한 CellType 라벨(Kupffer "
        "cells / Monocyte-derived macrophage / Neutrophils) 범위에서 골수계 세포만 골라내 "
        "작은 행렬로 추출했다.",
        result_fn=lambda s: add_result_console(s, [
            "Myeloid cells to extract:",
            "",
            "  Kupffer cells        15,842",
            "  Monocyte/MoMF         6,203",
            "  Neutrophils           1,153",
            "",
            "Extracted: ~32,000 genes x",
            "           23,198 myeloid cells",
            "-> gse166504_myeloid_counts.rds",
            "   gse166504_myeloid_meta.rds",
        ]),
        result_explain_kr="전체 데이터를 다 로드하지 않고 필요한 골수계 세포(23,198개)만 뽑아내, "
                           "이후 02단계에서 다루기 쉬운 크기로 재클러스터링할 수 있게 준비했다.",
    )

    emit_step(
        "ext/02_gse166504_cluster_annotate.R — Recluster GSE166504 myeloid compartment",
        os.path.join(EXT_SCRIPTS, "02_gse166504_cluster_annotate.R"),
        "추출된 23,198개 골수계 세포 범위에서 자체적으로 재클러스터링(resolution=0.6, 18개 "
        "클러스터)을 수행해, project1에서 정의한 S100A8hi/Neutrophil/Kupffer 구분 기준(동일한 "
        "마커 패널)이 저자의 원래 CellType 라벨과 독립적으로도 재현되는지 검증했다.",
        result_fn=lambda s: add_result_image(s, os.path.join(EXT_FIG, "gse166504_umap_authorlabel.png")),
        result_explain_kr="18개 클러스터와 저자의 3개 coarse 라벨을 cross-tab으로 대조 -- "
                           "S100a8 양성이 뚜렷한 클러스터 하나로 모이지 않고 여러 단핵구/대식세포 클러스터에 걸쳐 연속적 gradient로 나타남을 확인했다.",
    )

    emit_step(
        "ext/03_gse156057_qc_cluster.R — QC & cluster GSE156057 (CD45+-sorted)",
        os.path.join(EXT_SCRIPTS, "03_gse156057_qc_cluster.R"),
        "GSE156057(Remmerie et al. 2020 Immunity, CD45+ 정렬, SD/WD x 12/24/36주, 6개 샘플) "
        "h5 파일 6개 범위에서 QC(nFeature 200-6000, nCount<30000, percent.mt<25) 후 병합해 "
        "resolution=0.8로 40개 클러스터를 얻고, project1과 동일한 17개 유전자 마커 패널로 "
        "클러스터별 발현비율표를 만들었다.",
        result_fn=lambda s: add_result_image(s, os.path.join(EXT_FIG, "gse156057_umap_clusters.png")),
        result_explain_kr="CD45+ 정렬 덕분에 면역세포 해상도가 높아, 40개 클러스터 중 3개가 "
                           "S100A8hi 후보(S100a8/9 97-99%, Ly6g/Camp/Ltf/Ngp 모두 <5%) 기준에 명확히 부합했다.",
    )

    emit_step(
        "ext/04_gse156057_proportions.R — S100A8hi vs Kupffer vs Neutrophil proportions over time",
        os.path.join(EXT_SCRIPTS, "04_gse156057_proportions.R"),
        "03단계 pct_expr 표에서 S100a8/S100a9>0.9 & Ly6g/Camp/Ltf/Ngp<0.1(S100A8hi), "
        "S100a8/9>0.9 & Ly6g/Camp>0.5(Neutrophil), Clec4f>0.7 & S100a8<0.2(Kupffer) 조건 "
        "범위에서 각 클러스터를 분류하고, SD/WD x 12/24/36주 6개 샘플 전체의 세포class별 "
        "비율을 계산했다.",
        result_fn=lambda s: add_result_image(s, os.path.join(EXT_FIG, "gse156057_proportions.png")),
        result_explain_kr="S100A8hi_macrophage 비율이 WD에서 12주(1.45배), 24주(1.21배)엔 SD보다 "
                           "높았지만 36주(0.81배)엔 역전 -- diet-dependent 경향과 일치하되 단일 샘플이라 통계적 확정은 어려움.",
    )

    emit_step(
        "ext/05_summary_figure.R — Assemble the external-validation summary figure",
        os.path.join(EXT_SCRIPTS, "05_summary_figure.R"),
        "이미 계산된 CSV들(무거운 Seurat 객체를 다시 로드하지 않고) 범위에서 GSE156057의 마커 "
        "프로필(A)과 시간별 비율(B), GSE166504의 S100a8-vs-Csf1r gradient(C) 세 패널을 "
        "patchwork로 조합해 하나의 최종 그림으로 만들었다.",
        result_fn=lambda s: add_result_image(s, os.path.join(EXT_FIG, "Figure5_external_validation.png")),
        result_explain_kr="CD45+ 정렬 데이터(GSE156057)에서는 population이 재현되지만 whole-liver "
                           "데이터(GSE166504)에서는 이산적 클러스터가 아닌 연속적 gradient로만 나타나 -- 면역세포 정렬이 해상도에 핵심적임을 시사.",
    )

    # =============================================== PART 3: FOLLOW-UP ==
    section_divider("Part 3 — Follow-up: Origin, Trigger, Fate",
                     "원 논문이 명시적으로 미해결로 남긴 3가지 질문을 같은 데이터로 재검증")

    emit_step(
        "13_origin_signature_scoring.R — Origin: monocyte-recruitment vs. resident-Kupffer scoring",
        os.path.join(SCRIPTS, "13_origin_signature_scoring.R"),
        "문헌 기반(Guilliams 2022 Cell 등) monocyte-recruitment(Ly6c2/Ccr2/Sell/Plac8/Chil3) "
        "vs resident-Kupffer(Clec4f/Timd4/Vsig4/Cd5l/Marco) 마커 패널로 sanity check 하는 "
        "범위에서, 합쳐졌던 \"Macrophage_Kupffer\" 라벨의 4개 서브클러스터를 다시 쪼개 각각의 "
        "pct-expressing을 직접 비교했다.",
        result_fn=lambda s: add_result_image(s, os.path.join(FIG, "13_origin_marker_dotplot.png")),
        result_explain_kr="클러스터18(진짜 상주 Kupffer, Clec4f 99.7%/Ccr2 7%)과 클러스터5+11"
                           "(67%, Ccr2 74-78%인 단핵구유래/LAM-유사)로 나뉨을 발견 -- S100A8hi는 이 둘 중 어느 쪽 시그니처도 갖지 않았다.",
    )

    emit_step(
        "ext/06_origin_signature_scoring.R — Cross-dataset check of the origin finding",
        os.path.join(EXT_SCRIPTS, "06_origin_signature_scoring.R"),
        "13단계와 동일한 monocyte/resident 시그니처를 GSE156057(discrete 클러스터 3개 재확인)과 "
        "GSE166504(discrete 클러스터 없음, 클러스터 평균값 상관분석으로 대체) 범위에 각각 "
        "적용해, origin 발견이 다른 데이터셋에서도 일관되는지 검증했다.",
        result_fn=lambda s: add_result_image(s, os.path.join(EXT_FIG, "06_origin_scores_gse166504.png")),
        result_explain_kr="GSE166504에서는 resident score가 S100a8 gradient와 약한 음의 상관"
                           "(rho=-0.47, p=0.05)을 보여 부분적으로 일관된 방향을 시사했다.",
    )

    emit_step(
        "14_tf_activity_inference.R — Trigger: does C/EBP or AP-1 explain S100A8 induction?",
        os.path.join(SCRIPTS, "14_tf_activity_inference.R"),
        "원 논문이 제안만 하고 테스트하지 않은 C/EBP·AP-1 가설을 검증하기 위해 decoupleR+"
        "CollecTRI/DoRothEA 정식 TF activity inference를 시도했으나, OmnipathR의 organism "
        "데이터베이스 파싱 버그(GitHub 공식 이슈 #128/#121, 최신 개발버전 설치 후에도 재현)로 "
        "실패했다. 대안으로 C/EBP(Cebpa/b/d), AP-1(Fos/Fosb/Jun/Junb/Jund/Atf3) 유전자 자체의 "
        "mRNA 발현을 34개 bulk 샘플과 scRNA 범위에서 S100a9/S100A8hi 존재비와 상관분석했다.",
        result_fn=lambda s: add_result_image(s, os.path.join(FIG, "14_tf_bulk_correlation.png")),
        result_explain_kr="9개 TF 유전자 중 BH q<0.05에 도달한 것은 없었다 -- Cebpa는 오히려 "
                           "S100A8hi에서 가장 낮아(4.4%) \"C/EBP가 열쇠\"라는 가설과 정반대 방향이었다.",
    )

    emit_step(
        "15_fate_trajectory.R — Fate: transient recruitment or a maturing state?",
        os.path.join(SCRIPTS, "15_fate_trajectory.R"),
        "S100A8hi가 시간이 지나며 Kupffer/LAM 상태로 성숙하는지 확인하기 위해 두 가지를 했다: "
        "(a) GSE156057의 12/24/36주 시계열 범위에서 S100A8hi/Kupffer 비율의 시간 추세를 "
        "선형회귀로, (b) 정제된 4개 origin group만으로 새 PCA를 만들어 각 세포의 k=30 최근접"
        "이웃 구성을 origin group별로 집계하는 범위에서 전사체적 연속성을 검정했다.",
        result_fn=lambda s: add_result_image(s, os.path.join(FIG, "15_fate_umap_origin.png")),
        result_explain_kr="S100A8hi는 이웃의 96.2%가 자기 그룹(고립된 독립 상태)인 반면, "
                           "Kupffer_resident와 MoMF_LAM은 7.3%/2.6% 서로 섞여 -- 알려진 단핵구→Kupffer 분화 연속체와 일치. "
                           "S100A8hi는 이 궤적 위의 한 지점이 아니다.",
    )

    # ------------------------------------------------------------- summary
    summary_slide("Summary and Conclusions", [
        ("Primary manuscript finding", [
            "지방생성 유전자와 상관된 것은 S100A8hi 서브셋이 아니라 일반 대식세포/Kupffer 총량 (r=0.74 vs r=0.11)",
            "외부 검증: CD45+ 정렬 데이터에서는 재현, whole-liver 데이터에서는 gradient로만 관찰",
            "MKO axis2의 null 결과는 \"효과 없음\"이 아니라 \"underpowered\" (Cohen's d>=2.0-2.4만 검출 가능)",
        ]),
        ("Follow-up finding (origin / trigger / fate)", [
            "\"Macrophage_Kupffer\" 라벨이 실제로는 진짜 상주 Kupffer + Ccr2-high 단핵구유래/LAM 혼합이었음을 발견",
            "S100A8hi는 고전적 Ccr2-의존 모집 시그니처도, C/EBP·AP-1 유발도, 단핵구→Kupffer 분화 궤적도 아닌 독립 상태",
            "결론: 원 논문이 \"아마 최근 모집된 단핵구겠지\"로 넘긴 이 population은 아직 이름 붙지 않은 급성 반응 상태에 가깝다",
        ]),
    ])

    closing_slide("Thank You", "Questions & Discussion  |  bioinform25/project1")


if __name__ == "__main__":
    build()
    prs.save(OUT_PATH)
    print(f"Saved to {OUT_PATH}  ({len(prs.slides)} slides)")
