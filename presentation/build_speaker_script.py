#!/usr/bin/env python3
"""Injects a full speaker script (Korean, spoken register) into lab_meeting_full.pptx
as PowerPoint notes (visible in Presenter View), and also renders a standalone,
printable script document (speaker_script.docx) with slide number + title + script,
so the user can rehearse offline."""

import os
from pptx import Presentation
from docx import Document
from docx.shared import Pt, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH

HERE = os.path.dirname(__file__)
PPTX_PATH = os.path.join(HERE, "lab_meeting_full.pptx")
DOCX_PATH = os.path.join(HERE, "speaker_script.docx")

# One entry per slide, in exact deck order (51 slides).
SCRIPT = [
# 1. Title
"안녕하세요, 오늘은 최근에 진행한 S100A8+ 대식세포와 MASLD 관련 후속 분석 결과를 발표하겠습니다. "
"원 논문은 JCI에 실린 Guan 등의 연구고요, 저희는 이 논문이 공개한 데이터를 가지고 완전히 새로운 각도의 "
"분석을 해봤습니다. 오늘 발표는 크게 세 부분입니다 — 원래 저희가 했던 주 분석, 외부 데이터로 재검증한 부분, "
"그리고 최근에 진행한 후속 분석까지 코드 단위로 하나씩 다 보여드릴게요.",

# 2. Overall Study Flow
"먼저 전체 그림부터 말씀드릴게요. 원 논문의 핵심 스토리는 이겁니다 — 고지방식이로 지방세포가 죽으면, "
"그 신호로 간에 S100A8이라는 단백질을 발현하는 대식세포가 모이고, 이 세포들이 CCN3이라는 유전자를 억제해서 "
"결국 CD36이 올라가고 간에 지방이 쌓인다는 겁니다. 근데 이 논문은 CD36, CCN3 이렇게 개별 유전자 경로만 "
"bulk 데이터로 추적했어요. 저희는 여기서 질문을 하나 던졌습니다 — 이 마우스들에서 bulk RNA-seq이랑 "
"single-cell RNA-seq을 같이 가지고 있으니까, 유전자 하나가 아니라 '간의 전체 세포 구성 자체'가 "
"genotype에 따라 바뀌는지 볼 수 있지 않을까 하는 거였죠. 그래서 1부에서는 scRNA로 세포 시그니처를 만들어서 "
"bulk 34개 샘플에 역투영하는 deconvolution을 했고, 2부에서는 이 방법으로 찾은 population이 다른 독립된 "
"데이터셋에서도 재현되는지 검증했고, 3부에서는 원 논문이 스스로 '모르겠다'고 인정한 질문들을 저희가 "
"직접 파고들었습니다. 결론부터 살짝 말씀드리면, 지방생성 유전자와 진짜 연관된 건 S100A8 서브셋이 아니라 "
"일반 대식세포 전체였다는, 예상을 뒤집는 결과가 나왔습니다.",

# 3. Divider: Part 1
"자, 그럼 1부부터 시작하겠습니다. 원 논문이 쓴 데이터 그대로, GSE285614 bulk RNA-seq 34개 샘플이랑 "
"GSE285615 scRNA-seq, WT랑 TG 각 1개 샘플을 썼습니다.",

# 4. 01_scRNA_qc.R (1/2)
"먼저 scRNA 데이터부터 정리했습니다. 이 스크립트는 GEO에서 받은 h5 파일 두 개를 불러와서 Seurat 객체로 "
"만드는 부분이에요. 코드가 좀 길어서 두 장으로 나눴는데, 이 앞부분은 파일을 읽고 QC 지표를 계산하는 부분입니다.",

# 5. 01_scRNA_qc.R (2/2)
"뒷부분에서 실제로 필터링을 합니다. 세포당 유전자 수, 전체 read 수, 미토콘드리아 비율 이 세 가지 기준으로 "
"저품질 세포를 걸러냈어요. 오른쪽 violin plot을 보시면 TG랑 WT 두 샘플 다 정상적인 분포를 보이는 걸 "
"확인할 수 있습니다. 이 단계에서는 별다른 이상은 없었고, 그냥 표준적인 QC 절차라고 생각하시면 됩니다.",

# 6. 02_scRNA_cluster_annotate.R (1/2)
"다음은 클러스터링입니다. 샘플이 딱 WT 1개, TG 1개뿐이라서 사실 샘플이랑 genotype이 완전히 같은 거예요. "
"그래서 Harmony 같은 배치보정을 쓰면 저희가 보려는 진짜 생물학적 신호까지 지워버릴 위험이 있어서, "
"일부러 배치보정 없이 그냥 단순하게 병합만 하고 클러스터링을 돌렸습니다.",

# 7. 02_scRNA_cluster_annotate.R (2/2)
"결과적으로 28개 클러스터가 나왔고, 오른쪽 UMAP에서 보시는 것처럼 꽤 깔끔하게 나뉘었습니다. 간세포, "
"Kupffer세포, 내피세포 같은 주요 계통 마커 22개로 대분류를 확인했고, 클러스터별 top5 마커 유전자도 "
"뽑아서 다음 단계 라벨링 근거로 저장해뒀습니다.",

# 8. 03_scRNA_annotate_finalize.R (1/2)
"이 스크립트가 저희 분석에서 제일 중요한 단계 중 하나예요. 여기서 28개 클러스터를 실제 세포타입 12개로 "
"매핑합니다. 코드에서 보시는 것처럼 클러스터 번호 하나하나에 사람이 직접 마커 유전자를 대조해서 라벨을 붙였어요.",

# 9. 03_scRNA_annotate_finalize.R (2/2)
"여기서 핵심 포인트입니다 — S100a8, S100a9를 발현하는 세포가 사실 두 개의 완전히 다른 세포로 갈립니다. "
"하나는 저희가 S100A8hi_macrophage라고 부르는, 단핵구·대식세포 계통 세포고요, 다른 하나는 그냥 성숙한 "
"호중구예요. 오른쪽 dot plot 보시면 호중구만 Ly6g, Camp, Ltf, Ngp에 양성인 걸 보실 수 있는데, 이 둘을 "
"구분하지 않고 그냥 'S100A8 양성 세포'로 뭉뚱그리면 저희가 보려는 population이 호중구로 오염돼 버립니다. "
"그래서 이 구분이 이후 모든 분석의 기초가 됐어요.",

# 10. 04_scRNA_ground_truth_props.R
"다음은 scRNA에서 직접 관찰한 세포타입 구성비를 계산하는 단계입니다. WT, TG 각 샘플에서 세포타입별 "
"비율을 세서, 이걸 나중에 bulk deconvolution 결과랑 비교할 '정답지'로 쓰려고 했어요. 결과를 보면 "
"S100A8hi_macrophage가 TG에서 WT보다 1.46배 많게 나왔습니다. 근데 이건 샘플이 genotype당 딱 1개씩이라, "
"이 숫자를 얼마나 믿을 수 있는지는 나중에 다시 짚고 넘어가게 됩니다.",

# 11. 05_bulk_load_merge.R (1/2)
"이제 bulk 데이터로 넘어갑니다. GSE285614는 사실 두 개의 별도 xlsx 파일로 나뉘어 있어요 — 하나는 WT vs "
"Bcl2TG 비교고, 다른 하나는 WT vs MKO(대식세포 특이적 S100a8 결손) 비교입니다. 이 코드는 두 파일에서 "
"저자가 같이 올려둔 시료별 raw count를 뽑아내는 부분이에요.",

# 12. 05_bulk_load_merge.R (2/2)
"뒷부분에서 두 파일의 유전자셋을 교집합해서 하나의 count matrix로 합칩니다. 오른쪽 표가 최종 34개 "
"샘플의 구성인데요, 여기서 꼭 짚고 싶은 게 있습니다 — 두 파일 다 'WT' 그룹이 있는데, 이 둘은 서로 다른 "
"시퀀싱 배치라서 절대로 같은 시료로 취급하면 안 됩니다. axis라는 컬럼으로 항상 분리해서 다뤘어요.",

# 13. 06_deconvolution.R (1/2)
"드디어 핵심 분석입니다, deconvolution이요. scRNA 참조가 genotype당 딱 1개체뿐이라서, 개체 간 분산을 "
"추정해야 하는 MuSiC 같은 방법은 통계적으로 못 씁니다. 그래서 세포타입별 마커 유전자만 있으면 되는 "
"BisqueRNA의 marker-based 방법을 선택했어요. 이 앞부분은 FindAllMarkers로 세포타입별 마커를 뽑는 부분입니다.",

# 14. 06_deconvolution.R (2/2)
"뒷부분에서 실제로 34개 bulk 샘플에 이 마커 시그니처를 투영해서 상대적 세포타입 존재비 점수를 계산합니다. "
"오른쪽 표에서 보시는 것처럼 세포타입마다 마커 유전자 수가 5개에서 25개까지 다양한데, 이건 저희가 임의로 "
"정한 게 아니라 알고리즘이 PC1 분산을 최대화하는 방향으로 자동으로 고른 겁니다. 다만 여기서 나온 점수는 "
"세포타입 간 절대적인 비율이 아니라 상대적인 scale이라, 같은 세포타입을 여러 시료 간에 비교하는 데만 "
"써야 한다는 점이 중요합니다.",

# 15. 07_validate_axis1.R (1/3)
"이제 이 deconvolution 결과를 믿어도 되는지 검증하는 단계입니다. 사실 여기서 좀 흥미로운 일이 있었어요 — "
"04단계 scRNA ground truth는 TG에서 S100A8hi가 더 많다고 나왔는데, 06단계 deconvolution은 정반대로 WT가 "
"더 많다고 나왔거든요. 둘 중 뭐가 맞는지 확인해야 했습니다.",

# 16. 07_validate_axis1.R (2/3)
"그래서 저자가 원래 논문에서 정식으로 돌린 DESeq2 결과, 그러니까 n=5 대 n=5로 제대로 통계 검정을 한 "
"고신뢰 데이터를 직접 열어봤습니다. S100a8, Cd68, Cd36 같은 핵심 마커 유전자들의 방향을 하나하나 대조하는 "
"부분이 이 코드입니다.",

# 17. 07_validate_axis1.R (3/3)
"결과를 보시면, 원저자 DESeq2는 S100a8, Cd68, Cd36 전부 WT가 TG보다 높은 방향으로 나왔어요. 이건 저희 "
"06단계 deconvolution이랑 정확히 일치합니다. 그러니까 n=1짜리 scRNA 쌍이 개체차 때문에 우연히 반대 "
"방향으로 나온 거고, deconvolution 쪽이 맞다고 결론 내렸습니다. 이 검증 과정이 있었기 때문에 다음 "
"단계에서 deconvolution을 자신 있게 새로운 질문에 적용할 수 있었어요.",

# 18. 08_apply_axis2_MKO.R (1/2)
"검증이 끝났으니 이제 진짜 새로운 질문에 적용합니다. MKO, 그러니까 대식세포 특이적으로 S100a8을 없앤 "
"마우스는 scRNA 자체가 없어요. 그래서 방금 검증한 방법을 그대로 이 axis에 적용하는 게 이 코드입니다.",

# 19. 08_apply_axis2_MKO.R (2/2)
"결과는— S100A8hi_macrophage든 Macrophage_Kupffer든, WT랑 MKO 사이에 유의한 차이가 없었습니다. "
"처음엔 좀 실망스러운 결과였는데, 저희는 이걸 그냥 숨기지 않고 정직하게 null 결과로 보고했어요. "
"그리고 이게 정말 '효과가 없는 건지' 아니면 '표본이 작아서 못 보는 건지'는 뒤에 11단계에서 다시 짚습니다.",

# 20. 09_correlate_lipid_genes.R (1/2)
"이제 원 논문의 핵심 기전이 저희 deconvolution 결과와도 연결되는지 봤습니다. S100A8+ 대식세포가 결국 "
"CD36을 올려서 지방을 쌓는다는 게 원 논문 얘기니까, 저희 세포타입 존재비 점수가 Cd36 같은 지방생성 "
"유전자 발현이랑 상관관계가 있는지 34개 bulk 샘플 전체에서 계산했습니다.",

# 21. 09_correlate_lipid_genes.R (2/2)
"그런데 여기서 예상치 못한 결과가 나왔어요. S100A8hi_macrophage 존재비는 Cd36이랑 거의 상관이 없었습니다, "
"r값이 0.11 정도로요. 이게 사실 오늘 발표의 가장 중요한 반전의 시작점이에요 — 나중에 나오지만, "
"S100A8hi 서브셋이 아니라 일반 Macrophage_Kupffer 전체가 지방생성 유전자랑 강하게 연관돼 있었습니다.",

# 22. 11_statistical_rigor_addendum.R (1/3)
"원고를 쓰기 전에 저희 스스로 통계적으로 허술한 부분이 없는지 다시 점검했습니다. 이 스크립트가 그 재점검 "
"과정이에요. 먼저 09단계 지질유전자 상관분석은 30개나 되는 검정을 했는데 다중비교 보정이 빠져 있었어서, "
"그 부분을 BH-FDR로 보정하는 코드입니다.",

# 23. 11_statistical_rigor_addendum.R (2/3)
"그다음으로 07, 08단계처럼 표본이 4~5개밖에 안 되는 소표본 비교에는 정규성 가정이 깨지기 쉬우니까, "
"비모수 검정인 Wilcoxon을 1차 지표로 삼고 Cohen's d 효과크기랑 95% 신뢰구간을 같이 병기하도록 다시 "
"계산했습니다.",

# 24. 11_statistical_rigor_addendum.R (3/3)
"제일 중요한 부분이 이겁니다 — 08단계의 null 결과가 '진짜 효과가 없는 건지' 아니면 '표본이 작아서 "
"못 보는 건지' 구분하기 위해서, 80% 검정력 기준으로 이 표본크기에서 최소 검출 가능한 효과크기를 "
"계산했어요. 오른쪽 결과를 보시면 n=4, n=5에서 Cohen's d가 2.0에서 2.4는 돼야 검출이 가능합니다. "
"이건 굉장히 큰 효과크기거든요. 그래서 08단계 null 결과는 '효과 없음'이 아니라 '이 정도 표본으로는 "
"큰 효과 말고는 못 본다'라고 정직하게 재해석했습니다. 이게 저희 논문의 통계적 엄밀성을 보여주는 "
"중요한 대목이에요.",

# 25. Divider: Part 2
"1부는 여기까지고요, 2부로 넘어가겠습니다. 여기서는 저희가 찾은 S100A8hi_macrophage라는 population이 "
"정말 진짜인지, 아니면 이 마우스 한 쌍에서만 우연히 나온 artifact인지 독립된 외부 데이터 두 개로 "
"검증했습니다.",

# 26. ext/01_gse166504_extract_myeloid.R
"첫 번째 외부 데이터셋은 GSE166504인데, 진짜 생물학적 반복이 있는 마우스 간 데이터예요, 그룹당 3마리에서 "
"8마리까지요. 근데 전체 데이터가 너무 커서, 저자가 제공한 세포타입 라벨을 이용해서 골수계 세포만 먼저 "
"추출했습니다. Kupffer세포, 단핵구유래대식세포, 호중구 이렇게 세 종류만 골라서 23,198개 세포로 줄였어요.",

# 27. ext/02_gse166504_cluster_annotate.R (1/2)
"추출한 골수계 세포를 가지고 독자적으로 다시 클러스터링을 했습니다. 이 코드가 그 과정이에요 — 저희가 "
"1부에서 썼던 것과 똑같은 마커 패널로, 이번엔 완전히 다른 마우스들에서도 같은 구분이 재현되는지 보는 겁니다.",

# 28. ext/02_gse166504_cluster_annotate.R (2/2)
"결과가 좀 흥미로웠습니다. 18개 클러스터가 나왔는데, S100a8 양성 세포가 하나의 뚜렷한 클러스터로 모이지 "
"않고 여러 단핵구·대식세포 클러스터에 걸쳐서 연속적인 gradient로 나타났어요. 즉 이 population이 "
"'있긴 있는데' 이 데이터에서는 이산적인 클러스터로 딱 떨어지지 않았다는 겁니다. 이건 나중에 왜 그런지 "
"설명이 나옵니다.",

# 29. ext/03_gse156057_qc_cluster.R (1/2)
"두 번째 외부 데이터셋은 GSE156057인데, 이건 CD45라는 면역세포 마커로 정렬(sorting)까지 한 데이터예요. "
"식이랑 시점별로 6개 샘플이 있습니다. 이 코드는 6개 h5 파일을 QC하고 병합하는 부분이에요.",

# 30. ext/03_gse156057_qc_cluster.R (2/2)
"클러스터링 결과 40개나 되는 클러스터가 나왔는데, 면역세포로 미리 정렬을 해놔서 그런지 해상도가 훨씬 "
"좋았습니다. 저희가 쓰는 마커 패널을 적용해보니까 이 중 3개 클러스터가 S100A8hi 후보 조건에 명확하게 "
"들어맞았어요 — S100a8/9는 97~99% 양성인데 호중구 마커는 다 5% 미만인 클러스터들이요.",

# 31. ext/04_gse156057_proportions.R (1/2)
"자 그럼 이 population이 식이나 시점에 따라 어떻게 변하는지 봤습니다. 앞서 정의한 기준으로 "
"S100A8hi, Kupffer, 호중구를 분류하고, 6개 샘플 전체에서 비율을 계산하는 코드입니다.",

# 32. ext/04_gse156057_proportions.R (2/2)
"결과를 보시면, 서양식이(WD) 그룹에서 S100A8hi 비율이 12주, 24주까지는 표준식이보다 높았다가 "
"36주에서는 오히려 역전됩니다. 샘플이 시점당 1개뿐이라 통계적으로 확정할 순 없지만, 방향성은 "
"식이-의존적인 급성 반응이라는 저희 가설과 일치했어요.",

# 33. ext/05_summary_figure.R (1/2)
"이제 두 외부 데이터셋 결과를 하나의 그림으로 합칩니다. 무거운 Seurat 객체를 다시 불러오지 않고, "
"이미 계산해둔 CSV들만 가지고 그림을 그리는 효율적인 방식이에요.",

# 34. ext/05_summary_figure.R (2/2)
"최종 그림이 이겁니다. 왼쪽이 CD45+ 정렬된 GSE156057에서의 마커 프로필, 가운데가 시간에 따른 비율 변화, "
"오른쪽이 GSE166504에서 S100a8이 이산적 클러스터가 아니라 연속적 gradient로 나타난다는 걸 보여줍니다. "
"결론적으로 — 면역세포로 미리 정렬한 데이터에서는 이 population이 재현되고, whole-liver 데이터에서는 "
"gradient로만 보인다는 겁니다. 그러니까 이 population이 가짜가 아니라, 정렬(sorting)이 이걸 잡아내는 데 "
"핵심적인 요인이라는 뜻이에요.",

# 35. Divider: Part 3
"자, 여기까지가 원래 저희 논문 내용이고요. 이제 최근에 진행한 후속 분석으로 넘어가겠습니다. 원 논문의 "
"Discussion을 다시 꼼꼼히 읽어보니까, 저자들이 스스로 '이건 아직 모른다'고 인정한 부분이 세 가지 "
"있더라고요 — 이 세포가 어디서 왔는지(기원), 뭐가 이 세포를 유도하는지(유발인자), 그리고 이 세포가 "
"시간이 지나면 어떻게 되는지(운명). 저희가 이미 가진 데이터로 이 세 가지를 직접 파봤습니다.",

# 36. 13_origin_signature_scoring.R (1/3)
"첫 번째, 기원 문제입니다. 이 population이 순환하는 단핵구가 새로 모집된 건지, 아니면 원래 있던 세포가 "
"국소적으로 재프로그래밍된 건지 궁금했어요. 그래서 문헌에서 확립된 단핵구-모집 마커랑 상주-Kupffer 마커로 "
"점수를 매기는 코드를 짰습니다.",

# 37. 13_origin_signature_scoring.R (2/3)
"근데 이걸 돌리다가 sanity check 과정에서 예상 밖의 걸 발견했어요. 합쳐서 하나의 라벨로 취급했던 "
"'Macrophage_Kupffer'가 사실 단일 집단이 아니었습니다. 그래서 이 부분에서 원래 4개였던 서브클러스터를 "
"다시 쪼개서 각각 확인하는 코드로 발전시켰습니다.",

# 38. 13_origin_signature_scoring.R (3/3)
"이게 결과인데요, 정말 흥미롭습니다. 원래 하나로 뭉쳐 있던 Macrophage_Kupffer 중에서, 클러스터18은 "
"진짜 상주 Kupffer세포였고 — Clec4f 99.7%, Ccr2는 7%밖에 안 됐어요 — 클러스터5랑 11, 그러니까 전체의 "
"67%는 오히려 Ccr2가 74~78%나 되는, 단핵구유래로 보이는 세포들이었습니다. 그리고 정작 저희가 관심 있던 "
"S100A8hi_macrophage는 이 둘 중 어느 쪽 시그니처도 갖고 있지 않았어요. 이건 저희가 이미 발행한 논문의 "
"Discussion 문장 하나를 실제로 고치게 만든 발견이었습니다.",

# 39. ext/06_origin_signature_scoring.R (1/4)
"이 기원 관련 발견이 다른 데이터셋에서도 일관되는지 검증했습니다. 이 코드는 같은 시그니처를 GSE156057과 "
"GSE166504 두 데이터셋에 각각 적용하는 부분이에요.",

# 40. ext/06_origin_signature_scoring.R (2/4)
"GSE156057 쪽 코드가 이어집니다. 여기서도 Kupffer 라벨을 세분화해서 같은 패턴이 재현되는지 봤습니다.",

# 41. ext/06_origin_signature_scoring.R (3/4)
"그리고 GSE166504는 discrete 클러스터가 없으니까, 클러스터 평균값끼리 상관분석을 하는 방식으로 "
"바꿔서 적용한 부분이 여기예요.",

# 42. ext/06_origin_signature_scoring.R (4/4)
"결과를 보면, GSE166504에서는 상주-Kupffer 점수가 S100a8 gradient랑 약한 음의 상관(rho=-0.47, p=0.05)을 "
"보였어요. 완벽하게 깔끔하진 않지만, 방향은 저희 1차 발견이랑 부분적으로 일치했습니다.",

# 43. 14_tf_activity_inference.R (1/4)
"두 번째 질문, 유발인자입니다. 원 논문 Discussion에 'C/EBP나 AP-1이라는 전사인자를 테스트해보면 "
"흥미로울 것'이라고 써 있는데, 실제로 테스트는 안 했더라고요. 그래서 저희가 직접 해봤습니다. 원래는 "
"decoupleR이라는 패키지로 정식 전사인자 활성도 추론을 하려고 했는데, 이 부분이 그 시도입니다.",

# 44. 14_tf_activity_inference.R (2/4)
"근데 이게 OmnipathR이라는 의존 패키지의 버그로 계속 실패했어요. GitHub 최신 개발버전까지 다 설치해봤는데도 "
"안 되길래, 공식 이슈 트래커를 확인해보니까 이게 저희만의 문제가 아니라 이미 보고된 미해결 버그였습니다. "
"그래서 대안을 찾은 부분이 여기예요.",

# 45. 14_tf_activity_inference.R (3/4)
"대안으로 전사인자 자체의 mRNA 발현을 대리지표로 썼습니다. C/EBP 계열이랑 AP-1 계열 유전자들의 발현을 "
"bulk 샘플이랑 scRNA 양쪽에서 확인하는 코드입니다.",

# 46. 14_tf_activity_inference.R (4/4)
"결과는 이렇습니다 — 9개 전사인자 유전자 중에 통계적으로 유의한 게 하나도 없었어요. 심지어 Cebpa는, "
"문헌상 대식세포 염증 유도에 제일 중요하다는 그 유전자인데, 오히려 S100A8hi에서 가장 낮게 나왔습니다 "
"(4.4%). 'C/EBP가 열쇠일 것'이라는 가설과는 정반대 방향이었던 거죠. 이것도 정직하게 null 결과로 "
"보고했습니다.",

# 47. 15_fate_trajectory.R (1/3)
"마지막 질문, 운명입니다. 이 세포가 시간이 지나면서 Kupffer나 LAM 같은 성숙한 상태로 변하는 건지, "
"아니면 독립적인 상태로 계속 남아있는 건지 봤어요. 첫 번째 방법은 GSE156057의 시계열 데이터로 "
"시간에 따른 비율 변화 추세를 보는 겁니다.",

# 48. 15_fate_trajectory.R (2/3)
"두 번째 방법이 좀 더 정교한데요, 4개의 정제된 그룹만으로 새로 PCA를 만들어서, 각 세포의 최근접이웃 "
"30개가 어느 그룹에 속하는지 계산하는 코드입니다. 만약 S100A8hi가 다른 그룹으로 변해가는 중이라면, "
"이웃 중에 다른 그룹 세포가 꽤 섞여 있어야겠죠.",

# 49. 15_fate_trajectory.R (3/3)
"결과가 아주 명확했습니다. S100A8hi_macrophage는 이웃의 96.2%가 자기 자신과 같은 그룹이었어요 — 완전히 "
"고립된 독립적인 상태라는 뜻입니다. 반면에 진짜 상주 Kupffer랑 단핵구유래 그룹은 서로 7.3%, 2.6% 정도 "
"섞여 있었는데, 이건 문헌에 이미 알려진 단핵구가 Kupffer로 분화하는 연속적인 과정과 정확히 일치합니다. "
"그러니까 S100A8hi는 이 익숙한 분화 경로 위에 있는 한 지점이 아니라, 완전히 별개의 상태라는 거죠.",

# 50. Summary and Conclusions
"정리하겠습니다. 원래 저희 논문의 핵심 발견은, 지방생성 유전자와 상관된 게 S100A8hi 서브셋이 아니라 "
"일반 대식세포·Kupffer 총량이었다는 겁니다. 외부 데이터로도 이 population 자체는 재현이 됐고요, "
"MKO 그룹의 null 결과는 효과가 없는 게 아니라 표본크기 문제라는 것도 통계적으로 확인했습니다. "
"그리고 이번 후속 분석에서는, 그동안 하나로 뭉뚱그렸던 'Macrophage_Kupffer' 라벨 안에 사실 진짜 "
"상주세포랑 단핵구유래 세포가 섞여 있었다는 걸 발견했고, S100A8hi_macrophage는 고전적인 단핵구 모집 "
"경로도, C/EBP·AP-1 유도 경로도, 익숙한 분화 궤적도 아닌, 아직 이름이 붙지 않은 독립적인 급성 반응 "
"상태에 가깝다는 결론을 얻었습니다. 발표는 여기까지입니다, 질문 있으시면 편하게 해주세요.",

# 51. Thank You
"감사합니다.",
]


def inject_notes():
    prs = Presentation(PPTX_PATH)
    slides = list(prs.slides)
    assert len(slides) == len(SCRIPT), f"slide count {len(slides)} != script count {len(SCRIPT)}"
    for slide, text in zip(slides, SCRIPT):
        slide.notes_slide.notes_text_frame.text = text
    try:
        prs.save(PPTX_PATH)
        print(f"Injected speaker notes into {len(slides)} slides of {PPTX_PATH}")
    except PermissionError:
        tmp_path = PPTX_PATH.replace(".pptx", "_withnotes.pptx")
        prs.save(tmp_path)
        print(f"PERMISSION DENIED on {PPTX_PATH} (likely open in PowerPoint). "
              f"Saved to {tmp_path} instead -- close PowerPoint and re-run, "
              f"or manually replace the original with this file.")


def build_script_doc():
    prs = Presentation(PPTX_PATH)
    titles = []
    for s in prs.slides:
        t = ""
        for shp in s.shapes:
            if shp.has_text_frame and shp.text_frame.text.strip():
                t = shp.text_frame.text.strip().splitlines()[0]
                break
        titles.append(t)

    doc = Document()
    normal = doc.styles["Normal"]
    normal.font.name = "Malgun Gothic"
    normal.font.size = Pt(11)
    normal.paragraph_format.space_after = Pt(6)
    normal.paragraph_format.line_spacing = 1.35

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = title.add_run("발표 대본 — S100A8+ Macrophage & MASLD")
    r.bold = True
    r.font.size = Pt(20)
    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = sub.add_run("lab_meeting_full.pptx (51 slides) 대응 발표 스크립트")
    r2.italic = True
    r2.font.size = Pt(12)
    doc.add_paragraph()

    for i, (title_text, script_text) in enumerate(zip(titles, SCRIPT), 1):
        h = doc.add_paragraph()
        r = h.add_run(f"Slide {i} — {title_text}")
        r.bold = True
        r.font.size = Pt(13)
        h.paragraph_format.space_before = Pt(14)
        h.paragraph_format.space_after = Pt(4)
        p = doc.add_paragraph()
        p.add_run(script_text)
        p.paragraph_format.left_indent = Inches(0.2)

    doc.save(DOCX_PATH)
    print(f"Saved {DOCX_PATH}")


if __name__ == "__main__":
    inject_notes()
    build_script_doc()
