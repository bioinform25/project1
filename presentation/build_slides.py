#!/usr/bin/env python3
"""Builds lab_meeting.pptx: the narrative deck for tomorrow's lab meeting --
source paper -> motivation -> method -> main-manuscript results -> follow-up
(origin/trigger/fate) results -> synthesis. Built with python-pptx (no Node/
LibreOffice on this machine)."""

import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

REPO = os.path.join(os.path.dirname(__file__), "..")
FIG = os.path.join(REPO, "figures")
EXT_FIG = os.path.join(REPO, "external_validation", "figures")
OUT_PATH = os.path.join(os.path.dirname(__file__), "lab_meeting.pptx")

# ---------------------------------------------------------------- palette ---
NAVY = RGBColor(0x1E, 0x27, 0x61)
NAVY_DARK = RGBColor(0x14, 0x1A, 0x42)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CORAL = RGBColor(0xC1, 0x12, 0x1F)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
MUTED = RGBColor(0x6B, 0x72, 0x80)
CARD_BG = RGBColor(0xF3, 0xF6, 0xFC)
CARD_BG2 = RGBColor(0xFD, 0xEE, 0xEC)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

HEAD_FONT = "Cambria"
BODY_FONT = "Calibri"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def new_slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return s


def add_rect(slide, left, top, width, height, fill_color, radius=False, shadow=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    shp.shadow.inherit = False
    if radius:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    return shp


def add_text(slide, text, left, top, width, height, size, color, bold=False,
             align=PP_ALIGN.LEFT, font=BODY_FONT, italic=False, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.0):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        r.font.color.rgb = color
    return tb


def add_bullets(slide, items, left, top, width, height, size, color,
                 font=BODY_FONT, space_after=8, bullet_color=None, line_spacing=1.1):
    """items: list of (text, level) or plain strings (level 0)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        text, level = item if isinstance(item, tuple) else (item, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        indent = Inches(0.28 * level)
        p.alignment = PP_ALIGN.LEFT
        prefix = "•  " if level == 0 else "–  "
        r = p.add_run()
        r.text = prefix + text
        r.font.size = Pt(size - level * 1)
        r.font.name = font
        r.font.color.rgb = color
        pPr = p._pPr
        if pPr is None:
            pPr = p.get_or_add_pPr()
        pPr.set("marL", str(indent))
        pPr.set("indent", str(-Inches(0.28)))
    return tb


def add_image_fit(slide, path, left, top, max_w, max_h, align="center"):
    with Image.open(path) as im:
        iw, ih = im.size
    ar = iw / ih
    box_ar = max_w / max_h
    if ar > box_ar:
        w = max_w
        h = int(max_w / ar)
    else:
        h = max_h
        w = int(max_h * ar)
    x = left + (max_w - w) // 2 if align == "center" else left
    y = top + (max_h - h) // 2
    slide.shapes.add_picture(path, x, y, width=w, height=h)
    return x, y, w, h


def slide_number(slide, n, dark=False):
    color = ICE if dark else MUTED
    add_text(slide, str(n), Inches(12.7), Inches(7.1), Inches(0.5), Inches(0.3),
              11, color, align=PP_ALIGN.RIGHT)


def section_tag(slide, text, dark=False):
    color = ICE if dark else CORAL
    add_text(slide, text.upper(), Inches(0.6), Inches(0.42), Inches(8), Inches(0.35),
              13, color, bold=True, font=BODY_FONT)


def slide_title(slide, text, dark=False, top=Inches(0.75), size=30):
    color = WHITE if dark else NAVY
    add_text(slide, text, Inches(0.6), top, Inches(12.2), Inches(1.1), size, color,
              bold=True, font=HEAD_FONT)


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ======================================================== SLIDE 1 -- TITLE ==
def build_title():
    s = new_slide(NAVY)
    add_rect(s, Inches(0), Inches(5.55), SLIDE_W, Inches(1.95), NAVY_DARK)
    add_text(s, "S100A8+ 대식세포와 MASLD", Inches(0.9), Inches(2.15), Inches(11.5), Inches(1.0),
              40, WHITE, bold=True, font=HEAD_FONT)
    add_text(s, "scRNA 기반 Deconvolution으로 다시 본 지방세포 사멸-간 지방축적 축",
              Inches(0.9), Inches(3.05), Inches(11.5), Inches(0.7), 20, ICE, font=HEAD_FONT, italic=True)
    add_text(s, "원 논문: Guan et al., \"Adipocyte death promotes hepatic infiltration of "
                 "S100A8+ macrophages and steatotic liver disease progression in mice\", "
                 "J Clin Invest 2025;135(21):e190635",
              Inches(0.9), Inches(5.85), Inches(11.5), Inches(0.6), 13, ICE, font=BODY_FONT)
    add_text(s, "Lab Meeting  |  bioinform25/project1", Inches(0.9), Inches(6.35), Inches(11.5), Inches(0.5),
              13, RGBColor(0x8A, 0x93, 0xC4), font=BODY_FONT)
    return s


# ================================================= SLIDE 2 -- SOURCE PAPER ==
def build_paper_background():
    s = new_slide(WHITE)
    section_tag(s, "Background")
    slide_title(s, "원 논문이 밝힌 것: 지방세포 사멸 → S100A8+ 대식세포 → 간 지방축적")
    steps = [
        ("1", "지방세포 사멸", "HFD로 부고환 지방세포가 죽으면서\nFFA + 세포외소포(EV) 방출"),
        ("2", "S100A8+ 대식세포", "FFA/EV가 대식세포·단핵구에서\nS100A8 발현을 유도, 간에 축적"),
        ("3", "CCN3 억제", "S100A8+ 대식세포가 CCN3(CD36의\n음성조절인자)을 억제"),
        ("4", "CD36 상승 → 지방축적", "간세포 CD36 발현 상승 →\n지방산 흡수 증가 → MASLD 진행"),
    ]
    card_w = Inches(2.85)
    gap = Inches(0.22)
    x0 = Inches(0.6)
    y0 = Inches(2.35)
    card_h = Inches(3.55)
    for i, (num, head, desc) in enumerate(steps):
        x = x0 + i * (card_w + gap)
        add_rect(s, x, y0, card_w, card_h, CARD_BG, radius=True)
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.22), y0 + Inches(0.28), Inches(0.6), Inches(0.6))
        circ.fill.solid(); circ.fill.fore_color.rgb = CORAL; circ.line.fill.background(); circ.shadow.inherit = False
        tf = circ.text_frame; tf.margin_left = 0; tf.margin_right = 0
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num; r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = HEAD_FONT
        add_text(s, head, x + Inches(0.22), y0 + Inches(1.05), card_w - Inches(0.44), Inches(0.75),
                  16, NAVY, bold=True, font=HEAD_FONT)
        add_text(s, desc, x + Inches(0.22), y0 + Inches(1.85), card_w - Inches(0.44), Inches(1.55),
                  12.5, DARK_TEXT, font=BODY_FONT, line_spacing=1.25)
        if i < 3:
            arrow = s.shapes.add_shape(MSO_SHAPE.CHEVRON, x + card_w + Inches(0.01), y0 + Inches(1.55),
                                         Inches(0.2), Inches(0.45))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = MUTED; arrow.line.fill.background(); arrow.shadow.inherit = False
    add_text(s, "macrophage-specific S100a8 KO(MKO)는 CD36 발현·간 지방축적·혈청 ALT를 낮춤 -- "
                 "이 축을 유전학적으로도 검증",
              Inches(0.6), Inches(6.15), Inches(12.1), Inches(0.6), 13, MUTED, italic=True)
    slide_number(s, 2)
    add_notes(s, "원 논문의 핵심 기전을 4단계로 요약. 이후 슬라이드에서 우리가 이 중 어느 지점에 새로운 질문을 던졌는지 이어집니다.")
    return s


# =================================================== SLIDE 3 -- DATASETS ==
def build_datasets():
    s = new_slide(WHITE)
    section_tag(s, "Background")
    slide_title(s, "사용한 공개 데이터: GSE285614 (bulk) + GSE285615 (scRNA)")
    add_rect(s, Inches(0.6), Inches(2.3), Inches(5.8), Inches(4.3), CARD_BG, radius=True)
    add_text(s, "GSE285614  ·  Bulk RNA-seq", Inches(0.95), Inches(2.55), Inches(5.1), Inches(0.5),
              18, NAVY, bold=True, font=HEAD_FONT)
    add_bullets(s, [
        "34개 샘플, DESeq2 raw count 포함",
        ("Axis 1: WT vs Bcl2TG (지방세포사멸 저항), chow 3v3 / HFD 5v5", 1),
        ("Axis 2: WT vs 대식세포특이 S100a8 KO(MKO), chow 4v5 / HFD 5v4", 1),
        "두 axis는 서로 다른 시퀀싱 배치 -- 절대 풀링하지 않음",
    ], Inches(0.95), Inches(3.15), Inches(5.1), Inches(3.2), 14, DARK_TEXT, space_after=10)

    add_rect(s, Inches(6.75), Inches(2.3), Inches(5.95), Inches(4.3), CARD_BG2, radius=True)
    add_text(s, "GSE285615  ·  scRNA-seq", Inches(7.1), Inches(2.55), Inches(5.3), Inches(0.5),
              18, CORAL, bold=True, font=HEAD_FONT)
    add_bullets(s, [
        "WT_HFD vs TG_HFD 단 2개 샘플 (genotype당 n=1)",
        ("HFD 4개월, 10x 3' v3.1, 17,031 cells, 14개 세포타입", 1),
        "chow 조건 scRNA는 없음 -- axis-1(HFD)만 대응",
        "MKO axis는 scRNA 자체가 아예 존재하지 않음",
    ], Inches(7.1), Inches(3.15), Inches(5.3), Inches(3.2), 14, DARK_TEXT, space_after=10)
    add_text(s, "핵심 제약: scRNA 참조가 genotype당 n=1 -- subject-variance 기반 방법(MuSiC 등) 대신 "
                 "marker-gene 기반 deconvolution(BisqueRNA)을 선택한 이유",
              Inches(0.6), Inches(6.75), Inches(12.1), Inches(0.5), 13, MUTED, italic=True)
    slide_number(s, 3)
    return s


# ================================================= SLIDE 4 -- MOTIVATION ==
def build_motivation():
    s = new_slide(NAVY)
    section_tag(s, "Our Question", dark=True)
    slide_title(s, "원 논문이 안 던진 질문: \"세포 구성 자체가 바뀌는가?\"", dark=True)
    add_text(s,
        "원 논문은 CD36·CCN3라는 개별 유전자 경로만 bulk DEG로 추적했다. "
        "하지만 이 연구는 bulk와 scRNA를 같은 마우스 코호트에서 함께 갖고 있다 -- "
        "그렇다면 scRNA 참조로 만든 세포타입 시그니처를 bulk 34개 샘플에 역투영해서, "
        "\"간의 전체 세포구성 자체가 genotype에 따라 바뀌는가\"를 직접 물을 수 있지 않을까?",
        Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.7), 17, ICE, font=BODY_FONT, line_spacing=1.35)

    add_rect(s, Inches(0.9), Inches(4.25), Inches(5.6), Inches(2.4), NAVY_DARK, radius=True)
    add_text(s, "1단계 · 검증", Inches(1.2), Inches(4.5), Inches(5), Inches(0.4), 15, CORAL, bold=True, font=HEAD_FONT)
    add_text(s, "scRNA가 있는 Axis 1(WT vs Bcl2TG)에서\ndeconvolution 결과가 실제 scRNA "
                 "ground truth와 같은 방향인지 검증",
              Inches(1.2), Inches(5.0), Inches(5.0), Inches(1.5), 14, WHITE, line_spacing=1.3)

    add_rect(s, Inches(6.8), Inches(4.25), Inches(5.6), Inches(2.4), NAVY_DARK, radius=True)
    add_text(s, "2단계 · 적용", Inches(7.1), Inches(4.5), Inches(5), Inches(0.4), 15, CORAL, bold=True, font=HEAD_FONT)
    add_text(s, "검증된 방법을 scRNA가 없는 Axis 2(MKO)에\n적용해 새로운 질문에 답 -- "
                 "S100a8 손실이 세포구성을 바꾸는가?",
              Inches(7.1), Inches(5.0), Inches(5.0), Inches(1.5), 14, WHITE, line_spacing=1.3)
    slide_number(s, 4, dark=True)
    return s


# =================================================== SLIDE 5 -- METHOD ==
def build_method():
    s = new_slide(WHITE)
    section_tag(s, "Method")
    slide_title(s, "방법: Marker-Based Deconvolution (BisqueRNA)")
    add_bullets(s, [
        "scRNA 참조에서 세포타입별 마커 시그니처 행렬 구축",
        "MuSiC / Bisque ReferenceBasedDecomposition 등 subject-variance 기반 방법은 "
        "genotype당 n=1이라 통계적으로 부적절 -> 배제",
        "BisqueRNA::MarkerBasedDecomposition 채택: cross-subject 분산 없이 안정적 "
        "세포타입 시그니처만 있으면 적용 가능",
        "34개 bulk 샘플 전체에 시그니처를 투영해 상대적 세포타입 존재비 점수 산출",
        "이 한계(n=1 scRNA 참조)는 논문에 명시적으로 기술",
    ], Inches(0.6), Inches(2.5), Inches(7.6), Inches(4.3), 16, DARK_TEXT, space_after=16)

    add_rect(s, Inches(8.55), Inches(2.5), Inches(4.15), Inches(4.3), CARD_BG, radius=True)
    add_text(s, "왜 이 방법인가", Inches(8.85), Inches(2.75), Inches(3.6), Inches(0.4),
              15, NAVY, bold=True, font=HEAD_FONT)
    add_bullets(s, [
        "필요: 세포타입별 마커 유전자 시그니처",
        "불필요: subject 간 분산 추정",
        "결과: 상대적 존재비 점수 (절대 세포수 아님)",
    ], Inches(8.85), Inches(3.3), Inches(3.6), Inches(3.2), 13, DARK_TEXT, space_after=12)
    slide_number(s, 5)
    return s


# ============================================ SLIDE 6 -- FIGURE 1 (scRNA) ==
def build_figure1():
    s = new_slide(WHITE)
    section_tag(s, "Results · Figure 1")
    slide_title(s, "S100A8hi_macrophage는 호중구와 완전히 다른 세포다", size=26)
    add_image_fit(s, os.path.join(FIG, "Figure1.png"), Inches(0.6), Inches(2.0), Inches(7.6), Inches(4.9))
    add_rect(s, Inches(8.55), Inches(2.0), Inches(4.15), Inches(4.9), CARD_BG, radius=True)
    add_bullets(s, [
        "18,167개 세포, 12개 세포타입으로 주석",
        "S100a8/S100a9 양성 세포가 두 개의 완전히 "
        "다른 계통으로 분리됨:",
        ("S100A8hi_macrophage (Ly6g 음성, "
         "단핵구/대식세포 계통)", 1),
        ("Neutrophil (Ly6g/Camp/Ltf/Ngp 양성, "
         "성숙 호중구)", 1),
        "이 둘을 구분하지 않으면 핵심 population이 "
        "호중구에 오염됨",
    ], Inches(8.85), Inches(2.3), Inches(3.6), Inches(4.4), 13.5, DARK_TEXT, space_after=12)
    slide_number(s, 6)
    return s


# ======================================== SLIDE 7 -- FIGURE 2 (axis-1) ==
def build_figure2():
    s = new_slide(WHITE)
    section_tag(s, "Results · Figure 2")
    slide_title(s, "1단계 검증: Deconvolution이 실제 scRNA와 같은 방향", size=26)
    add_image_fit(s, os.path.join(FIG, "Figure2.png"), Inches(0.6), Inches(2.0), Inches(7.6), Inches(4.9))
    add_rect(s, Inches(8.55), Inches(2.0), Inches(4.15), Inches(4.9), CARD_BG, radius=True)
    add_bullets(s, [
        "Axis 1 (WT vs Bcl2TG, HFD, n=5/그룹)",
        "Macrophage_Kupffer 존재비: WT > TG "
        "(Wilcoxon p=0.012, BH q=0.024, Cohen's d=-2.86)",
        "이는 원 논문의 replicated bulk DESeq2 결과와 "
        "방향이 정확히 일치",
        "-> 검증 통과, Axis 2(MKO)에 적용해도 된다는 근거 확보",
    ], Inches(8.85), Inches(2.3), Inches(3.6), Inches(4.4), 13.5, DARK_TEXT, space_after=12)
    slide_number(s, 7)
    return s


# ======================================== SLIDE 8 -- FIGURE 3 (axis-2) ==
def build_figure3():
    s = new_slide(WHITE)
    section_tag(s, "Results · Figure 3")
    slide_title(s, "2단계 적용: MKO는 유의한 세포구성 변화를 보이지 않음", size=25)
    add_image_fit(s, os.path.join(FIG, "Figure3.png"), Inches(0.6), Inches(2.0), Inches(7.6), Inches(4.9))
    add_rect(s, Inches(8.55), Inches(2.0), Inches(4.15), Inches(4.9), CARD_BG2, radius=True)
    add_bullets(s, [
        "S100A8hi_macrophage, Macrophage_Kupffer 모두 "
        "WT vs MKO 유의차 없음 (BH q>0.4)",
        "정직하게 보고: \"효과 없음\"이 아니라 "
        "\"이 표본크기(n=4-5)로는 못 봄\"",
        "사후 power 분석: 이 코호트는 Cohen's d ≥ 2.0-2.4 "
        "정도만 검출 가능",
        "-> null 결과를 과소적합력(underpowered) "
        "문제로 정확히 재해석",
    ], Inches(8.85), Inches(2.3), Inches(3.6), Inches(4.4), 13.5, DARK_TEXT, space_after=12)
    slide_number(s, 8)
    return s


# ==================================== SLIDE 9 -- FIGURE 4 (HEADLINE) ==
def build_figure4():
    s = new_slide(NAVY)
    section_tag(s, "Results · Figure 4 — 핵심 발견", dark=True)
    slide_title(s, "일반 대식세포 존재비가, S100A8hi 서브셋이 아니라, 지방생성 유전자와 상관", dark=True, size=23)
    add_image_fit(s, os.path.join(FIG, "Figure4.png"), Inches(0.6), Inches(2.15), Inches(7.4), Inches(4.75))
    add_rect(s, Inches(8.35), Inches(2.15), Inches(4.35), Inches(4.75), NAVY_DARK, radius=True)
    add_bullets(s, [
        "34개 bulk 중 19개 HFD 샘플 전체에서:",
        ("Macrophage_Kupffer vs Cd36: r=0.74, "
         "BH q=0.006", 1),
        ("Pparg r=0.75, Scd1 r=0.73, Plin2 r=0.65 "
         "(모두 q<0.05)", 1),
        ("S100A8hi_macrophage vs Cd36: r=0.11, "
         "q=0.76 (무관)", 1),
        "-> 이게 이 연구의 논문 제목이 된 반전 결과",
    ], Inches(8.65), Inches(2.45), Inches(3.85), Inches(4.3), 13.5, WHITE, space_after=12)
    slide_number(s, 9, dark=True)
    return s


# ====================================== SLIDE 10 -- FIGURE 5 (external) ==
def build_figure5():
    s = new_slide(WHITE)
    section_tag(s, "Results · Figure 5")
    slide_title(s, "외부 검증: CD45+ 정렬 데이터에선 재현, whole-liver에선 미재현", size=24)
    add_image_fit(s, os.path.join(FIG, "Figure5.png"), Inches(0.6), Inches(2.05), Inches(11.9), Inches(3.6))
    add_rect(s, Inches(0.6), Inches(5.85), Inches(11.9), Inches(1.15), CARD_BG, radius=True)
    add_text(s,
        "GSE156057(Remmerie, CD45+ 정렬, SD/WD x 12/24/36주)에서는 3개 클러스터로 재현, "
        "diet-dependent 경향. GSE166504(Su, whole-liver, 진짜 생물학적 반복 n=3-8)에서는 "
        "이산적 클러스터 대신 연속적 gradient만 관찰 -- 면역세포 정렬(sorting)이 이 population을 "
        "해상도 있게 잡아내는 데 핵심적임을 시사.",
        Inches(0.9), Inches(6.05), Inches(11.3), Inches(0.85), 13.5, DARK_TEXT, line_spacing=1.3)
    slide_number(s, 10)
    return s


# =================================== SLIDE 11 -- TITLE / SUMMARY ==
def build_summary1():
    s = new_slide(NAVY)
    section_tag(s, "Manuscript", dark=True)
    slide_title(s, "1부 결론 및 최종 제목", dark=True)
    add_rect(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.0), NAVY_DARK, radius=True)
    add_text(s, "\"Cross-Validated Deconvolution Identifies General Hepatic Macrophage "
                 "Burden, Not an S100A8-Hi Subset, as the Lipogenic-Gene Correlate in "
                 "Mouse MASLD\"",
              Inches(1.25), Inches(2.5), Inches(10.8), Inches(1.5), 20, WHITE, italic=True,
              font=HEAD_FONT, line_spacing=1.3)
    add_bullets(s, [
        "scRNA-informed deconvolution을 2단계(검증→적용)로 설계해 n=1 scRNA 참조의 한계를 "
        "정직하게 다룸",
        "핵심 발견: 지방생성 유전자와 상관된 것은 S100A8hi 서브셋이 아니라 일반 대식세포/"
        "Kupffer 총량 -- 예상을 뒤집는 결과",
        "MKO axis2의 null 결과를 power 분석으로 올바르게 재해석 (효과 없음 ≠ 검출력 부족)",
        "독립된 외부 데이터셋 2개로 S100A8hi population 자체의 재현성 검증",
    ], Inches(0.9), Inches(4.55), Inches(11.5), Inches(2.6), 15, ICE, space_after=12)
    slide_number(s, 11, dark=True)
    return s


# ============================== SLIDE 12 -- FOLLOW-UP MOTIVATION ==
def build_followup_motivation():
    s = new_slide(WHITE)
    section_tag(s, "Part 2 · Follow-up")
    slide_title(s, "여기서 더 파본다: 원 논문이 \"모른다\"고 인정한 3가지")
    add_text(s,
        "원 논문의 Discussion을 직접 읽고, 저자들이 명시적으로 미해결로 남긴 질문 3가지를 "
        "골라 같은 GEO 데이터(새 다운로드 없이)로 파고들었다.",
        Inches(0.6), Inches(2.0), Inches(12.1), Inches(0.7), 16, DARK_TEXT, line_spacing=1.3)
    cards = [
        ("기원 (Origin)", "S100A8+ 대식세포가 순환 단핵구 모집인지, 이주 전 단핵구 국소유도인지, "
                          "상주 대식세포 재프로그래밍인지 \"detailed investigations... will be essential\""),
        ("유발인자 (Trigger)", "\"It would be interesting to test whether... C/EBP and AP-1\" "
                              "-- 제안만 하고 실제로 테스트하지 않음"),
        ("운명 (Fate)", "HFD 4개월 단일 시점만 봐서, 이 population이 안정 상태인지 "
                        "다른 상태로 성숙하는 전이 상태인지 전혀 모름"),
    ]
    x0 = Inches(0.6); w = Inches(3.95); gap = Inches(0.2); y0 = Inches(3.0); h = Inches(3.7)
    colors = [CARD_BG, CARD_BG2, CARD_BG]
    for i, (head, desc) in enumerate(cards):
        x = x0 + i * (w + gap)
        add_rect(s, x, y0, w, h, colors[i], radius=True)
        add_text(s, head, x + Inches(0.3), y0 + Inches(0.35), w - Inches(0.6), Inches(0.6),
                  19, NAVY, bold=True, font=HEAD_FONT)
        add_text(s, desc, x + Inches(0.3), y0 + Inches(1.15), w - Inches(0.6), Inches(2.3),
                  13.5, DARK_TEXT, line_spacing=1.35)
    slide_number(s, 12)
    add_notes(s, "이 3가지는 GitHub 이슈처럼 '아직 안 풀린 질문'으로, 우리가 이미 가진 데이터와 파이프라인으로 직접 검증 가능했습니다.")
    return s


# ==================================== SLIDE 13 -- MODULE 1a (origin) ==
def build_module1a():
    s = new_slide(WHITE)
    section_tag(s, "Module 1 · Origin")
    slide_title(s, "\"Macrophage_Kupffer\"는 사실 하나의 세포집단이 아니었다", size=25)
    add_image_fit(s, os.path.join(FIG, "13_origin_marker_dotplot.png"), Inches(0.6), Inches(2.0), Inches(7.6), Inches(4.9))
    add_rect(s, Inches(8.55), Inches(2.0), Inches(4.15), Inches(4.9), CARD_BG2, radius=True)
    add_bullets(s, [
        "문헌기반 monocyte-recruitment(Ccr2 등) vs "
        "resident-Kupffer(Clec4f 등) 시그니처로 sanity "
        "check 하다가 발견",
        "합쳐졌던 \"Macrophage_Kupffer\" 4개 서브클러스터를 "
        "분리해보니:",
        ("클러스터 18: 진짜 상주 Kupffer (Clec4f 99.7%, "
         "Ccr2 7%)", 1),
        ("클러스터 5+11 (전체의 67%): Ccr2-high "
         "단핵구유래/LAM-유사", 1),
        "-> 원래 라벨을 Kupffer_resident vs "
        "Macrophage_MoMF_LAM으로 재분리",
    ], Inches(8.85), Inches(2.3), Inches(3.6), Inches(4.4), 13, DARK_TEXT, space_after=10)
    slide_number(s, 13)
    return s


def build_module1b():
    s = new_slide(WHITE)
    section_tag(s, "Module 1 · Origin")
    slide_title(s, "S100A8hi_macrophage는 고전적 Ccr2 모집 시그니처도 아니다", size=24)
    tbl_data = [
        ["", "S100A8hi", "Kupffer_resident", "MoMF_LAM", "Neutrophil"],
        ["Ccr2", "3.0%", "14.5%", "76.6%", "7.7%"],
        ["Clec4f", "17.8%", "70.3%", "20.6%", "11.0%"],
        ["Trem2 / Gpnmb", "0.9% / 2.1%", "20.6% / 10.1%", "52.5% / 30.0%", "2.2% / 2.2%"],
    ]
    left, top, w, h = Inches(0.6), Inches(2.15), Inches(7.6), Inches(2.6)
    gshape = s.shapes.add_table(len(tbl_data), len(tbl_data[0]), left, top, w, h)
    table = gshape.table
    table.columns[0].width = Inches(1.7)
    for c in range(1, 5):
        table.columns[c].width = Inches(1.475)
    for r_i, row in enumerate(tbl_data):
        for c_i, val in enumerate(row):
            cell = table.cell(r_i, c_i)
            cell.text = val if val else " "
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.runs[0]
            run.font.size = Pt(14)
            run.font.name = BODY_FONT
            run.font.bold = (r_i == 0)
            run.font.color.rgb = WHITE if r_i == 0 else DARK_TEXT
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if r_i == 0 else (CARD_BG if r_i % 2 == 0 else WHITE)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_rect(s, Inches(0.6), Inches(5.05), Inches(7.6), Inches(1.9), CARD_BG2, radius=True)
    add_text(s, "S100A8hi는 Ccr2가 가장 낮은 그룹 -- \"최근 모집된 단핵구\"라는 배제법 추론과 정반대. "
                 "LAM(Trem2/Gpnmb)도, 상주 Kupffer(Clec4f)도 아니고, 고전적 단핵구 모집 시그니처도 없음.",
              Inches(0.9), Inches(5.3), Inches(7.0), Inches(1.4), 14, DARK_TEXT, line_spacing=1.3)
    add_image_fit(s, os.path.join(FIG, "15_fate_umap_lamscore.png"), Inches(8.55), Inches(2.15), Inches(4.15), Inches(4.8))
    slide_number(s, 14)
    add_notes(s, "이 발견 덕분에 이미 push된 논문의 Discussion 한 문장(recently recruited monocyte-derived)도 실제로 고쳤습니다.")
    return s


# ==================================== SLIDE 15 -- MODULE 2 (trigger) ==
def build_module2():
    s = new_slide(WHITE)
    section_tag(s, "Module 2 · Trigger")
    slide_title(s, "C/EBP·AP-1이 S100A8 유도의 열쇠라는 가설, 지지 안 됨", size=25)
    add_image_fit(s, os.path.join(FIG, "14_tf_bulk_correlation.png"), Inches(0.6), Inches(2.0), Inches(7.6), Inches(4.9))
    add_rect(s, Inches(8.55), Inches(2.0), Inches(4.15), Inches(4.9), CARD_BG, radius=True)
    add_bullets(s, [
        "원 논문 Discussion: \"C/EBP나 AP-1을 테스트해보면 "
        "흥미로울 것\" (실제 테스트는 안 함)",
        "decoupleR + CollecTRI/DoRothEA 정식 TF activity "
        "inference 시도 -> OmnipathR 공식 미해결 버그로 실패 "
        "(GitHub 이슈 #128, #121 확인)",
        "대안: TF mRNA 발현을 직접 대리지표로 사용",
        "결과: bulk 34샘플 중 BH q<0.05 도달한 TF 없음",
        "scRNA: Cebpa가 오히려 S100A8hi에서 가장 낮음 "
        "(4.4%, MoMF_LAM은 51.1%) -- 가설과 정반대",
    ], Inches(8.85), Inches(2.3), Inches(3.6), Inches(4.4), 12.5, DARK_TEXT, space_after=9)
    slide_number(s, 15)
    add_notes(s, "패키지 버그(OmnipathR)는 GitHub 개발버전까지 다 설치해서 확인한 진짜 업스트림 이슈입니다 -- 로컬 문제가 아니었습니다.")
    return s


# ==================================== SLIDE 16 -- MODULE 3 (fate) ==
def build_module3():
    s = new_slide(WHITE)
    section_tag(s, "Module 3 · Fate")
    slide_title(s, "S100A8hi는 단핵구→Kupffer 분화 궤적과 분리된 독립 상태", size=23)
    add_image_fit(s, os.path.join(FIG, "15_fate_umap_origin.png"), Inches(0.6), Inches(2.0), Inches(5.7), Inches(4.9))
    tbl_data = [
        ["자기그룹", "S100A8hi", "Kupffer_res.", "MoMF_LAM"],
        ["S100A8hi_macrophage", "96.2%", "2.0%", "1.6%"],
        ["Kupffer_resident", "0.6%", "92.0%", "7.3%"],
        ["Macrophage_MoMF_LAM", "0.3%", "2.6%", "97.0%"],
    ]
    left, top, w, h = Inches(6.55), Inches(2.15), Inches(6.15), Inches(2.2)
    gshape = s.shapes.add_table(len(tbl_data), len(tbl_data[0]), left, top, w, h)
    table = gshape.table
    table.columns[0].width = Inches(2.35)
    for c in range(1, 4):
        table.columns[c].width = Inches(1.27)
    for r_i, row in enumerate(tbl_data):
        for c_i, val in enumerate(row):
            cell = table.cell(r_i, c_i)
            cell.text = val if val else " "
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            run = para.runs[0]
            run.font.size = Pt(12.5)
            run.font.name = BODY_FONT
            run.font.bold = (r_i == 0) or (c_i == 0)
            run.font.color.rgb = WHITE if r_i == 0 else DARK_TEXT
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if r_i == 0 else (CARD_BG if r_i % 2 == 0 else WHITE)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_rect(s, Inches(6.55), Inches(4.55), Inches(6.15), Inches(2.4), CARD_BG2, radius=True)
    add_text(s,
        "각 세포의 PCA 최근접이웃 30개 구성. S100A8hi는 96.2% 자기그룹 이웃(고립된 독립 상태). "
        "반면 Kupffer_resident ↔ MoMF_LAM은 7.3%/2.6% 서로 섞임 -- 이는 문헌상 알려진 단핵구→"
        "Kupffer 분화 연속체와 일치. S100A8hi는 이 궤적 위의 한 지점이 아니라는 뜻.",
        Inches(6.85), Inches(4.78), Inches(5.6), Inches(2.05), 13, DARK_TEXT, line_spacing=1.3)
    slide_number(s, 16)
    return s


# ============================================== SLIDE 17 -- SYNTHESIS ==
def build_synthesis():
    s = new_slide(NAVY)
    section_tag(s, "Synthesis", dark=True)
    slide_title(s, "세 모듈이 수렴하는 하나의 그림", dark=True)
    rows = [
        ("기원", "고전적 Ccr2-의존 단핵구 모집 시그니처가 없음"),
        ("유발인자", "C/EBP·AP-1 어느 쪽으로도 깔끔히 설명 안 됨"),
        ("운명", "단핵구→Kupffer/LAM 분화 궤적과 전사체적으로 분리됨"),
    ]
    y0 = Inches(2.35)
    for i, (head, desc) in enumerate(rows):
        y = y0 + i * Inches(1.35)
        add_rect(s, Inches(0.9), y, Inches(11.5), Inches(1.15), NAVY_DARK, radius=True)
        add_text(s, head, Inches(1.25), y + Inches(0.28), Inches(2.0), Inches(0.6),
                  17, CORAL, bold=True, font=HEAD_FONT)
        add_text(s, desc, Inches(3.5), y + Inches(0.3), Inches(8.5), Inches(0.6),
                  15.5, WHITE, font=BODY_FONT)
    add_text(s,
        "원 논문이 \"아마 최근 모집된 단핵구겠지\"라고 배제법으로 넘겼던 이 population은, "
        "실제로는 아직 이름 붙지 않은 독립적인 급성 반응 상태에 가깝다.",
        Inches(0.9), Inches(6.55), Inches(11.5), Inches(0.75), 16, ICE, italic=True, line_spacing=1.3)
    slide_number(s, 17, dark=True)
    return s


# ============================================ SLIDE 18 -- LIMITATIONS ==
def build_limitations():
    s = new_slide(WHITE)
    section_tag(s, "Limitations & Next Steps")
    slide_title(s, "한계와 다음 단계")
    add_rect(s, Inches(0.6), Inches(2.1), Inches(5.9), Inches(4.6), CARD_BG, radius=True)
    add_text(s, "한계", Inches(0.95), Inches(2.35), Inches(5.2), Inches(0.5), 17, NAVY, bold=True, font=HEAD_FONT)
    add_bullets(s, [
        "scRNA 참조가 genotype당 n=1 (사람이 아니라 이 프로젝트의 근본 제약)",
        "GSE156057 구성비 시계열은 diet당 timepoint 3개뿐 -- 방향성만, 통계 불가",
        "TF 활성도는 정식 activity inference가 아니라 mRNA 발현 대리지표",
        "MKO axis2는 여전히 underpowered (Cohen's d ≥ 2.0-2.4만 검출)",
    ], Inches(0.95), Inches(3.0), Inches(5.2), Inches(3.5), 14, DARK_TEXT, space_after=12)

    add_rect(s, Inches(6.85), Inches(2.1), Inches(5.9), Inches(4.6), CARD_BG2, radius=True)
    add_text(s, "다음 단계", Inches(7.2), Inches(2.35), Inches(5.2), Inches(0.5), 17, CORAL, bold=True, font=HEAD_FONT)
    add_bullets(s, [
        "OmnipathR 버그 패치되면 정식 TF activity inference 재시도",
        "더 큰 MKO 코호트로 재분석 (현재 검출력 한계 극복)",
        "S100A8hi_macrophage 고유 리간드-수용체 신호 (CellChat, 이미 설치됨) 탐색",
        "Neutrophil/Mast_cell 등 axis2에서 큰 효과크기 보인 세포타입 후속 추적",
    ], Inches(7.2), Inches(3.0), Inches(5.2), Inches(3.5), 14, DARK_TEXT, space_after=12)
    slide_number(s, 18)
    return s


# ================================================ SLIDE 19 -- THANK YOU ==
def build_thanks():
    s = new_slide(NAVY)
    add_text(s, "감사합니다", Inches(0.9), Inches(3.0), Inches(11.5), Inches(1.2),
              44, WHITE, bold=True, font=HEAD_FONT)
    add_text(s, "Questions & Discussion", Inches(0.9), Inches(4.0), Inches(11.5), Inches(0.6),
              20, ICE, font=HEAD_FONT, italic=True)
    add_text(s, "bioinform25/project1  |  전체 코드·데이터·수치는 manuscript.docx / supplemental_data.docx 참조",
              Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.5), 13, RGBColor(0x8A, 0x93, 0xC4))
    return s


# ================================================================= BUILD ==
build_title()
build_paper_background()
build_datasets()
build_motivation()
build_method()
build_figure1()
build_figure2()
build_figure3()
build_figure4()
build_figure5()
build_summary1()
build_followup_motivation()
build_module1a()
build_module1b()
build_module2()
build_module3()
build_synthesis()
build_limitations()
build_thanks()

prs.save(OUT_PATH)
print(f"Saved to {OUT_PATH}  ({len(prs.slides)} slides)")

